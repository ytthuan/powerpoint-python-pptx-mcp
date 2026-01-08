"""Tools for editing PPTX content."""

from pathlib import Path
from typing import Any, Dict

from mcp.types import Tool
from pptx import Presentation
from pptx.util import Inches, Pt

from ..core.pptx_handler import PPTXHandler
from ..utils.validators import validate_pptx_path, validate_slide_number, validate_position, validate_size, validate_image_path


def get_edit_tools() -> list[Tool]:
    """Get all edit tools."""
    return [
        Tool(
            name="update_slide_text",
            description="Update text in a specific shape on a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "pptx_path": {
                        "type": "string",
                        "description": "Path to the PPTX file",
                    },
                    "slide_number": {
                        "type": "integer",
                        "description": "Slide number (1-indexed)",
                        "minimum": 1,
                    },
                    "shape_id": {
                        "type": "integer",
                        "description": "Shape ID to update (from read_slide_content)",
                    },
                    "new_text": {
                        "type": "string",
                        "description": "New text content",
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Output PPTX path (optional, defaults to input path with .edited suffix)",
                    },
                },
                "required": ["pptx_path", "slide_number", "shape_id", "new_text"],
            },
        ),
        Tool(
            name="replace_slide_image",
            description="Replace an existing image on a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "pptx_path": {
                        "type": "string",
                        "description": "Path to the PPTX file",
                    },
                    "slide_number": {
                        "type": "integer",
                        "description": "Slide number (1-indexed)",
                        "minimum": 1,
                    },
                    "shape_id": {
                        "type": "integer",
                        "description": "Shape ID of the image to replace",
                    },
                    "new_image_path": {
                        "type": "string",
                        "description": "Path to the new image file",
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Output PPTX path (optional)",
                    },
                },
                "required": ["pptx_path", "slide_number", "shape_id", "new_image_path"],
            },
        ),
        Tool(
            name="add_text_box",
            description="Add a new text box to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "pptx_path": {
                        "type": "string",
                        "description": "Path to the PPTX file",
                    },
                    "slide_number": {
                        "type": "integer",
                        "description": "Slide number (1-indexed)",
                        "minimum": 1,
                    },
                    "text": {
                        "type": "string",
                        "description": "Text content",
                    },
                    "position": {
                        "type": "object",
                        "properties": {
                            "x": {"type": "number", "description": "X position in inches"},
                            "y": {"type": "number", "description": "Y position in inches"},
                        },
                    },
                    "size": {
                        "type": "object",
                        "properties": {
                            "width": {"type": "number", "description": "Width in inches"},
                            "height": {"type": "number", "description": "Height in inches"},
                        },
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Output PPTX path (optional)",
                    },
                },
                "required": ["pptx_path", "slide_number", "text"],
            },
        ),
        Tool(
            name="add_image",
            description="Add a new image to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "pptx_path": {
                        "type": "string",
                        "description": "Path to the PPTX file",
                    },
                    "slide_number": {
                        "type": "integer",
                        "description": "Slide number (1-indexed)",
                        "minimum": 1,
                    },
                    "image_path": {
                        "type": "string",
                        "description": "Path to the image file",
                    },
                    "position": {
                        "type": "object",
                        "properties": {
                            "x": {"type": "number", "description": "X position in inches"},
                            "y": {"type": "number", "description": "Y position in inches"},
                        },
                    },
                    "size": {
                        "type": "object",
                        "properties": {
                            "width": {"type": "number", "description": "Width in inches"},
                            "height": {"type": "number", "description": "Height in inches"},
                        },
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Output PPTX path (optional)",
                    },
                },
                "required": ["pptx_path", "slide_number", "image_path"],
            },
        ),
    ]


async def handle_update_slide_text(arguments: Dict[str, Any]) -> Dict[str, Any]:
    """Handle update_slide_text tool call."""
    pptx_path = validate_pptx_path(arguments["pptx_path"])
    slide_number = arguments["slide_number"]
    shape_id = arguments["shape_id"]
    new_text = arguments["new_text"]
    output_path = arguments.get("output_path")
    
    pres = Presentation(str(pptx_path))
    validate_slide_number(slide_number, len(pres.slides))
    
    slide = pres.slides[slide_number - 1]
    
    # Find shape by ID
    shape = None
    for s in slide.shapes:
        if s.shape_id == shape_id:
            shape = s
            break
    
    if not shape:
        raise ValueError(f"Shape with ID {shape_id} not found on slide {slide_number}")
    
    if not hasattr(shape, "text_frame") or not shape.text_frame:
        raise ValueError(f"Shape {shape_id} does not have a text frame")
    
    # Update text
    shape.text_frame.text = new_text
    
    # Save
    if output_path:
        output_path = Path(output_path)
    else:
        output_path = pptx_path.with_name(pptx_path.stem + ".edited.pptx")
    
    pres.save(str(output_path))
    
    return {
        "success": True,
        "output_path": str(output_path),
        "slide_number": slide_number,
        "shape_id": shape_id,
    }


async def handle_replace_slide_image(arguments: Dict[str, Any]) -> Dict[str, Any]:
    """Handle replace_slide_image tool call."""
    pptx_path = validate_pptx_path(arguments["pptx_path"])
    slide_number = arguments["slide_number"]
    shape_id = arguments["shape_id"]
    new_image_path = validate_image_path(arguments["new_image_path"])
    output_path = arguments.get("output_path")
    
    pres = Presentation(str(pptx_path))
    validate_slide_number(slide_number, len(pres.slides))
    
    slide = pres.slides[slide_number - 1]
    
    # Find shape by ID
    from pptx.shapes.picture import Picture
    shape = None
    for s in slide.shapes:
        if s.shape_id == shape_id:
            if isinstance(s, Picture):
                shape = s
                break
    
    if not shape:
        raise ValueError(f"Image shape with ID {shape_id} not found on slide {slide_number}")
    
    # Get position and size
    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height
    
    # Remove old shape
    sp = shape._element
    sp.getparent().remove(sp)
    
    # Add new image
    slide.shapes.add_picture(str(new_image_path), left, top, width, height)
    
    # Save
    if output_path:
        output_path = Path(output_path)
    else:
        output_path = pptx_path.with_name(pptx_path.stem + ".edited.pptx")
    
    pres.save(str(output_path))
    
    return {
        "success": True,
        "output_path": str(output_path),
        "slide_number": slide_number,
    }


async def handle_add_text_box(arguments: Dict[str, Any]) -> Dict[str, Any]:
    """Handle add_text_box tool call."""
    pptx_path = validate_pptx_path(arguments["pptx_path"])
    slide_number = arguments["slide_number"]
    text = arguments["text"]
    position = validate_position(arguments.get("position"))
    size = validate_size(arguments.get("size"))
    output_path = arguments.get("output_path")
    
    pres = Presentation(str(pptx_path))
    validate_slide_number(slide_number, len(pres.slides))
    
    slide = pres.slides[slide_number - 1]
    
    # Add text box
    left = Inches(position["x"])
    top = Inches(position["y"])
    width = Inches(size["width"])
    height = Inches(size["height"])
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = text
    
    # Save
    if output_path:
        output_path = Path(output_path)
    else:
        output_path = pptx_path.with_name(pptx_path.stem + ".edited.pptx")
    
    pres.save(str(output_path))
    
    return {
        "success": True,
        "output_path": str(output_path),
        "slide_number": slide_number,
        "shape_id": text_box.shape_id,
    }


async def handle_add_image(arguments: Dict[str, Any]) -> Dict[str, Any]:
    """Handle add_image tool call."""
    pptx_path = validate_pptx_path(arguments["pptx_path"])
    slide_number = arguments["slide_number"]
    image_path = validate_image_path(arguments["image_path"])
    position = validate_position(arguments.get("position"))
    size = validate_size(arguments.get("size"))
    output_path = arguments.get("output_path")
    
    pres = Presentation(str(pptx_path))
    validate_slide_number(slide_number, len(pres.slides))
    
    slide = pres.slides[slide_number - 1]
    
    # Add image
    left = Inches(position["x"])
    top = Inches(position["y"])
    width = Inches(size["width"])
    height = Inches(size["height"])
    
    picture = slide.shapes.add_picture(str(image_path), left, top, width, height)
    
    # Save
    if output_path:
        output_path = Path(output_path)
    else:
        output_path = pptx_path.with_name(pptx_path.stem + ".edited.pptx")
    
    pres.save(str(output_path))
    
    return {
        "success": True,
        "output_path": str(output_path),
        "slide_number": slide_number,
        "shape_id": picture.shape_id,
    }


