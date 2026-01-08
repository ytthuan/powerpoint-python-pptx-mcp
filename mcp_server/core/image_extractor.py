"""Image extraction utilities for PPTX files."""

import base64
import io
import zipfile
from pathlib import Path
from typing import Any, Dict, List, Optional

from PIL import Image

from ..utils.validators import validate_pptx_path, validate_slide_number


def extract_images_from_pptx(pptx_path: str | Path, slide_number: Optional[int] = None) -> List[Dict[str, Any]]:
    """Extract images from PPTX file."""
    pptx_path = validate_pptx_path(pptx_path)
    
    images = []
    with zipfile.ZipFile(pptx_path, "r") as zip_file:
        # List all image files in the PPTX
        image_extensions = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".wmf", ".emf"}
        
        for file_info in zip_file.infolist():
            file_path = Path(file_info.filename)
            if file_path.suffix.lower() in image_extensions:
                # Check if it's in media folder (embedded images)
                if "ppt/media/" in file_info.filename:
                    try:
                        image_data = zip_file.read(file_info.filename)
                        # Try to get image info
                        img = Image.open(io.BytesIO(image_data))
                        width, height = img.size
                        format_name = img.format or "UNKNOWN"
                        
                        images.append({
                            "path": file_info.filename,
                            "filename": file_path.name,
                            "size": len(image_data),
                            "width": width,
                            "height": height,
                            "format": format_name,
                        })
                    except Exception as e:
                        # Skip invalid images
                        continue
    
    return images


def extract_image_as_base64(pptx_path: str | Path, image_path_in_zip: str) -> Optional[str]:
    """Extract a specific image from PPTX and return as base64."""
    pptx_path = validate_pptx_path(pptx_path)
    
    with zipfile.ZipFile(pptx_path, "r") as zip_file:
        try:
            image_data = zip_file.read(image_path_in_zip)
            return base64.b64encode(image_data).decode("utf-8")
        except KeyError:
            return None


def extract_slide_images(pptx_path: str | Path, slide_number: int) -> List[Dict[str, Any]]:
    """Extract images associated with a specific slide."""
    from pptx import Presentation
    
    pptx_path = validate_pptx_path(pptx_path)
    pres = Presentation(str(pptx_path))
    validate_slide_number(slide_number, len(pres.slides))
    
    slide = pres.slides[slide_number - 1]
    images = []
    
    # Extract images from shapes
    for shape in slide.shapes:
        if hasattr(shape, "image"):
            try:
                image = shape.image
                # Get image blob
                image_blob = image.blob
                img = Image.open(io.BytesIO(image_blob))
                width, height = img.size
                
                images.append({
                    "shape_id": shape.shape_id,
                    "name": shape.name if hasattr(shape, "name") else "",
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                    "image_width": width,
                    "image_height": height,
                    "format": img.format or "UNKNOWN",
                })
            except Exception:
                continue
    
    return images

