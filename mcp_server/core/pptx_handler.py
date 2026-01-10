"""Core PPTX file operations handler."""

from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..cache import PresentationCache
from ..config import get_config
from ..utils.validators import validate_pptx_path, validate_slide_number


class PPTXHandler:
    """Handler for PPTX file operations with optional caching support.

    This handler provides core PPTX file operations with:
    - Lazy loading of presentations
    - Optional caching for better performance
    - Automatic cache invalidation on modifications
    """

    def __init__(
        self,
        pptx_path: str | Path,
        cache: Optional[PresentationCache] = None,
        enable_cache: Optional[bool] = None,
    ):
        """Initialize handler with PPTX file path.

        Args:
            pptx_path: Path to PPTX file
            cache: Optional PresentationCache instance
            enable_cache: Whether to enable caching (uses config default if None)
        """
        self.pptx_path = validate_pptx_path(pptx_path)
        self._presentation: Optional[Presentation] = None

        # Caching support
        config = get_config()
        self._enable_cache = (
            enable_cache if enable_cache is not None else config.performance.enable_cache
        )
        self._cache = cache if self._enable_cache else None
        self._is_modified = False

    @property
    def presentation(self) -> Presentation:
        """Lazy load presentation with optional caching."""
        if self._presentation is None:
            # Try to get from cache first
            if self._cache and self._enable_cache:
                cached = self._cache.get_presentation(self.pptx_path)
                if cached is not None:
                    self._presentation = cached
                    return self._presentation

            # Load from file
            self._presentation = Presentation(str(self.pptx_path))

            # Cache it
            if self._cache and self._enable_cache:
                self._cache.cache_presentation(self.pptx_path, self._presentation)

        return self._presentation

    def reload(self):
        """Reload presentation from file.

        Clears cached presentation and marks for reload on next access.
        """
        self._presentation = None
        self._is_modified = False

        # Invalidate cache
        if self._cache and self._enable_cache:
            self._cache.invalidate(self.pptx_path)

    def get_slide_count(self) -> int:
        """Get total number of slides."""
        return len(self.presentation.slides)

    def is_slide_hidden(self, slide_number: int) -> bool:
        """Check if a slide is hidden.

        Args:
            slide_number: Slide number (1-indexed)

        Returns:
            True if slide is hidden, False if visible
        """
        validate_slide_number(slide_number, self.get_slide_count())
        slide = self.presentation.slides[slide_number - 1]
        # The 'show' attribute is on the <sld> element, not on <sldId>
        # '0' means hidden, missing or '1' means visible
        # Handle both string and potential integer/None cases
        show_attr = slide.element.get("show")
        if show_attr is None:
            return False  # Missing attribute means visible
        # Convert to string and check - can be '0', 0, '1', 1, etc.
        show_value = str(show_attr).strip()
        # Check for '0' (hidden) - after str() conversion, 0 becomes '0'
        return show_value == "0"

    def set_slide_hidden(self, slide_number: int, hidden: bool):
        """Set slide visibility.

        Args:
            slide_number: Slide number (1-indexed)
            hidden: True to hide the slide, False to show it
        """
        validate_slide_number(slide_number, self.get_slide_count())
        slide = self.presentation.slides[slide_number - 1]
        # The 'show' attribute is on the <sld> element, not on <sldId>
        if hidden:
            slide.element.set("show", "0")
        else:
            # Remove the show attribute to restore default (visible)
            if "show" in slide.element.attrib:
                del slide.element.attrib["show"]

        # Mark as modified
        self._is_modified = True

    def get_presentation_info(self) -> Dict[str, Any]:
        """Get presentation metadata."""
        slide_count = self.get_slide_count()
        visible_count = 0
        for idx in range(slide_count):
            slide = self.presentation.slides[idx]
            # The 'show' attribute is on the <sld> element, not on <sldId>
            # '0' means hidden, missing or '1' means visible
            show_attr = slide.element.get("show")
            if show_attr is None:
                visible_count += 1  # Missing attribute means visible
            else:
                show_value = str(show_attr).strip()
                if show_value != "0":
                    visible_count += 1
        hidden_count = slide_count - visible_count

        return {
            "file_path": str(self.pptx_path),
            "file_name": self.pptx_path.name,
            "slide_count": slide_count,
            "visible_slides": visible_count,
            "hidden_slides": hidden_count,
            "slide_size": {
                "width": self.presentation.slide_width,
                "height": self.presentation.slide_height,
            },
        }

    def get_slide_text(self, slide_number: int) -> str:
        """Extract all text from a slide."""
        validate_slide_number(slide_number, self.get_slide_count())
        slide = self.presentation.slides[slide_number - 1]
        return self._extract_text_from_slide(slide)

    def _extract_text_from_slide(self, slide) -> str:
        """Recursively extract text from all shapes in a slide."""
        texts = []
        for shape in slide.shapes:
            texts.extend(self._extract_text_from_shape(shape))
        return "\n".join(texts)

    def _extract_text_from_shape(self, shape: BaseShape) -> List[str]:
        """Recursively extract text from a shape."""
        texts = []
        if hasattr(shape, "text_frame") and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs)
                if text.strip():
                    texts.append(text)
        elif isinstance(shape, GroupShape):
            for sub_shape in shape.shapes:
                texts.extend(self._extract_text_from_shape(sub_shape))
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for row in shape.table.rows:
                row_texts = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_texts.append(cell_text)
                if row_texts:
                    texts.append(" | ".join(row_texts))
        return texts

    def get_slide_content(self, slide_number: int) -> Dict[str, Any]:
        """Get comprehensive content from a slide."""
        validate_slide_number(slide_number, self.get_slide_count())
        slide = self.presentation.slides[slide_number - 1]

        # Extract title
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text

        # Extract text
        text = self._extract_text_from_slide(slide)

        # Extract shapes info
        shapes_info = []
        for idx, shape in enumerate(slide.shapes):
            shape_info = {
                "index": idx,
                "shape_id": shape.shape_id,
                "shape_type": str(shape.shape_type),
                "name": shape.name if hasattr(shape, "name") else "",
            }
            if hasattr(shape, "text_frame") and shape.text_frame:
                shape_info["text"] = "".join(
                    run.text for para in shape.text_frame.paragraphs for run in para.runs
                )
            if isinstance(shape, Picture):
                shape_info["image"] = True
                shape_info["image_path"] = (
                    shape.image.filename if hasattr(shape.image, "filename") else None
                )
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                shape_info["table"] = True
                shape_info["rows"] = len(shape.table.rows)
                shape_info["columns"] = len(shape.table.columns)
            shapes_info.append(shape_info)

        return {
            "slide_number": slide_number,
            "title": title,
            "text": text,
            "shapes": shapes_info,
            "hidden": self.is_slide_hidden(slide_number),
        }

    def get_slide_images(self, slide_number: int) -> List[Dict[str, Any]]:
        """Get image information from a slide."""
        validate_slide_number(slide_number, self.get_slide_count())
        slide = self.presentation.slides[slide_number - 1]

        images = []
        for idx, shape in enumerate(slide.shapes):
            if isinstance(shape, Picture):
                image_info = {
                    "index": idx,
                    "shape_id": shape.shape_id,
                    "name": shape.name if hasattr(shape, "name") else "",
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                }
                if hasattr(shape.image, "filename"):
                    image_info["filename"] = shape.image.filename
                images.append(image_info)
            elif isinstance(shape, GroupShape):
                # Check nested shapes in groups
                for sub_shape in shape.shapes:
                    if isinstance(sub_shape, Picture):
                        image_info = {
                            "index": idx,
                            "shape_id": sub_shape.shape_id,
                            "name": sub_shape.name if hasattr(sub_shape, "name") else "",
                            "left": sub_shape.left,
                            "top": sub_shape.top,
                            "width": sub_shape.width,
                            "height": sub_shape.height,
                        }
                        if hasattr(sub_shape.image, "filename"):
                            image_info["filename"] = sub_shape.image.filename
                        images.append(image_info)
        return images

    def get_notes(self, slide_number: Optional[int] = None) -> Dict[str, Any]:
        """Get notes for slide(s)."""
        if slide_number is not None:
            validate_slide_number(slide_number, self.get_slide_count())
            slide = self.presentation.slides[slide_number - 1]
            notes_text = ""
            if slide.has_notes_slide:
                try:
                    notes_text = slide.notes_slide.notes_text_frame.text
                except Exception:
                    pass
            return {
                "slide": slide_number,
                "notes": notes_text,
            }
        else:
            # Return all notes
            slides = []
            for idx, slide in enumerate(self.presentation.slides, start=1):
                notes_text = ""
                if slide.has_notes_slide:
                    try:
                        notes_text = slide.notes_slide.notes_text_frame.text
                    except Exception:
                        pass
                slides.append({"slide": idx, "notes": notes_text})
            return {"slides": slides}

    def get_slides_metadata(self, include_hidden: bool = True) -> List[Dict[str, Any]]:
        """Get metadata for all slides including their visibility status.

        Args:
            include_hidden: If False, only return visible slides

        Returns:
            List of slide metadata dictionaries
        """
        slides_metadata = []
        slide_count = self.get_slide_count()
        self.presentation.slides._sldIdLst

        for i in range(slide_count):
            # Check visibility directly without validation overhead
            slide = self.presentation.slides[i]
            # The 'show' attribute is on the <sld> element, not on <sldId>
            show_attr = slide.element.get("show")
            if show_attr is None:
                is_hidden = False  # Missing attribute means visible
            else:
                show_value = str(show_attr).strip()
                is_hidden = show_value == "0"

            if not include_hidden and is_hidden:
                continue

            slide = self.presentation.slides[i]
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text

            slides_metadata.append(
                {
                    "slide_number": i + 1,  # Convert to 1-indexed
                    "title": title,
                    "hidden": is_hidden,
                    "slide_id": slide.slide_id,
                }
            )

        return slides_metadata

    def save(self, output_path: Optional[str | Path] = None):
        """Save presentation to file.

        Args:
            output_path: Optional output path. If None, saves to original path.
        """
        save_path = Path(output_path) if output_path else self.pptx_path
        self.presentation.save(str(save_path))

        # Invalidate cache after save
        if self._cache and self._enable_cache:
            self._cache.invalidate(save_path)
            # If saving to original path or same path, also invalidate original reference
            if output_path is None or Path(output_path) == self.pptx_path:
                self._cache.invalidate(self.pptx_path)

        self._is_modified = False
