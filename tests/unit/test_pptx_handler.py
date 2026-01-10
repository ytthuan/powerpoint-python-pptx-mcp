import pytest
from pathlib import Path
from unittest.mock import MagicMock, patch
from src.mcp_server.core.pptx_handler import PPTXHandler
from pptx.enum.shapes import MSO_SHAPE_TYPE

@pytest.fixture
def mock_config():
    with patch("src.mcp_server.core.pptx_handler.get_config") as mock:
        config = MagicMock()
        config.performance.enable_cache = False
        mock.return_value = config
        yield config

@pytest.fixture
def mock_validate_path():
    with patch("src.mcp_server.core.pptx_handler.validate_pptx_path") as mock:
        mock.return_value = Path("test.pptx")
        yield mock

@patch("pptx.Presentation")
def test_handler_init(mock_pres, mock_validate_path, mock_config):
    handler = PPTXHandler("test.pptx")
    assert handler.pptx_path == Path("test.pptx")
    assert handler._presentation is None

@patch("src.mcp_server.core.pptx_handler.Presentation")
def test_handler_lazy_load(mock_pres_class, mock_validate_path, mock_config):
    handler = PPTXHandler("test.pptx")
    mock_pres = mock_pres_class.return_value
    
    # First access loads it
    pres = handler.presentation
    mock_pres_class.assert_called_once()
    assert pres == mock_pres
    
    # Second access uses cached version
    pres2 = handler.presentation
    mock_pres_class.assert_called_once()
    assert pres2 == mock_pres

@patch("src.mcp_server.core.pptx_handler.Presentation")
def test_get_slide_count(mock_pres_class, mock_validate_path, mock_config):
    handler = PPTXHandler("test.pptx")
    handler.presentation.slides = [MagicMock(), MagicMock()]
    assert handler.get_slide_count() == 2

@patch("src.mcp_server.core.pptx_handler.Presentation")
def test_is_slide_hidden(mock_pres_class, mock_validate_path, mock_config):
    handler = PPTXHandler("test.pptx")
    mock_slide = MagicMock()
    handler.presentation.slides = [mock_slide]
    
    # Hidden
    mock_slide.element.get.return_value = "0"
    assert handler.is_slide_hidden(1) is True
    
    # Visible (missing show attr)
    mock_slide.element.get.return_value = None
    assert handler.is_slide_hidden(1) is False
    
    # Visible (show="1")
    mock_slide.element.get.return_value = "1"
    assert handler.is_slide_hidden(1) is False

@patch("src.mcp_server.core.pptx_handler.Presentation")
def test_set_slide_hidden(mock_pres_class, mock_validate_path, mock_config):
    handler = PPTXHandler("test.pptx")
    mock_slide = MagicMock()
    handler.presentation.slides = [mock_slide]
    
    # Set to hidden
    handler.set_slide_hidden(1, True)
    mock_slide.element.set.assert_called_with("show", "0")
    assert handler._is_modified is True
    
    # Set to visible
    mock_slide.element.attrib = {"show": "0"}
    handler.set_slide_hidden(1, False)
    assert "show" not in mock_slide.element.attrib

@patch("src.mcp_server.core.pptx_handler.Presentation")
def test_get_slide_text(mock_pres_class, mock_validate_path, mock_config):
    handler = PPTXHandler("test.pptx")
    mock_slide = MagicMock()
    handler.presentation.slides = [mock_slide]
    
    mock_shape = MagicMock()
    mock_para = MagicMock()
    mock_run = MagicMock()
    mock_run.text = "Hello World"
    mock_para.runs = [mock_run]
    mock_shape.text_frame.paragraphs = [mock_para]
    mock_slide.shapes = [mock_shape]
    
    text = handler.get_slide_text(1)
    assert text == "Hello World"

@patch("src.mcp_server.core.pptx_handler.Presentation")
def test_get_notes(mock_pres_class, mock_validate_path, mock_config):
    handler = PPTXHandler("test.pptx")
    mock_slide = MagicMock()
    mock_slide.has_notes_slide = True
    mock_slide.notes_slide.notes_text_frame.text = "The note"
    handler.presentation.slides = [mock_slide]
    
    result = handler.get_notes(1)
    assert result == {"slide": 1, "notes": "The note"}

@patch("src.mcp_server.core.pptx_handler.Presentation")
def test_save(mock_pres_class, mock_validate_path, mock_config):
    handler = PPTXHandler("test.pptx")
    # Trigger lazy load
    mock_pres = handler.presentation
    
    handler.save("output.pptx")
    mock_pres.save.assert_called_with("output.pptx")
    assert handler._is_modified is False
