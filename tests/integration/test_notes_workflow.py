import pytest
import tempfile
import os
from pptx import Presentation
from src.mcp_server.tools.notes_tools import handle_process_notes_workflow
from src.mcp_server.core.pptx_handler import PPTXHandler


@pytest.fixture
def test_pptx():
    """Create a test PPTX with 2 slides and note placeholders."""
    prs = Presentation()

    # Slide 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    # Force creation of notes slide and its placeholder
    notes_slide1 = slide1.notes_slide
    notes_slide1.notes_text_frame.text = "Initial note 1"

    # Slide 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    notes_slide2 = slide2.notes_slide
    notes_slide2.notes_text_frame.text = "Initial note 2"

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        tmp_path = tmp.name

    yield tmp_path

    if os.path.exists(tmp_path):
        os.unlink(tmp_path)


@pytest.mark.asyncio
async def test_process_notes_workflow_vietnamese_policy(test_pptx):
    """Test process_notes_workflow with Vietnamese speaker notes policy."""

    # Vietnamese content following the policy:
    # 1. Address audience as "anh/chị"
    # 2. Use "chúng ta" for shared actions
    # 3. Two versions: Short version + Original

    notes_data = [
        {
            "slide_number": 1,
            "short_text": "Chào anh/chị, chúng ta sẽ bắt đầu buổi thuyết trình hôm nay.",
            "original_text": (
                "Chào mừng các anh/chị đã đến với buổi chia sẻ. "
                "Hôm nay chúng ta sẽ cùng nhau tìm hiểu về giải pháp mới này."
            ),
        }
    ]

    result = await handle_process_notes_workflow(
        {"pptx_path": test_pptx, "notes_data": notes_data, "in_place": True}
    )

    assert result["success"] is True
    assert result["workflow"] == "process_notes_workflow"
    assert result["formatted_slides"] == 1

    # Verify the saved content
    handler = PPTXHandler(test_pptx)
    notes = handler.get_notes(1)

    expected_structure = (
        "- Short version:\n"
        "Chào anh/chị, chúng ta sẽ bắt đầu buổi thuyết trình hôm nay.\n\n"
        "- Original:\n"
        "Chào mừng các anh/chị đã đến với buổi chia sẻ. "
        "Hôm nay chúng ta sẽ cùng nhau tìm hiểu về giải pháp mới này."
    )

    assert notes["notes"] == expected_structure
    assert "anh/chị" in notes["notes"]
    assert "chúng ta" in notes["notes"]
    assert "- Short version:" in notes["notes"]
    assert "- Original:" in notes["notes"]


@pytest.mark.asyncio
async def test_process_notes_workflow_multi_slide(test_pptx):
    """Test process_notes_workflow with multiple slides."""

    notes_data = [
        {"slide_number": 1, "short_text": "Short 1", "original_text": "Original 1"},
        {"slide_number": 2, "short_text": "Short 2", "original_text": "Original 2"},
    ]

    result = await handle_process_notes_workflow(
        {"pptx_path": test_pptx, "notes_data": notes_data, "in_place": True}
    )

    assert result["success"] is True
    assert result["updated_slides"] == 2

    handler = PPTXHandler(test_pptx)
    notes1 = handler.get_notes(1)
    notes2 = handler.get_notes(2)

    assert "Original 1" in notes1["notes"]
    assert "Original 2" in notes2["notes"]
