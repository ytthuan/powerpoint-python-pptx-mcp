"""Tests for slide visibility features."""

import tempfile
import os
from pptx import Presentation

from mcp_server.core.pptx_handler import PPTXHandler


import pytest


@pytest.fixture
def temp_pptx():
    """Fixture to create a temporary PPTX file."""
    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    yield path
    if os.path.exists(path):
        os.unlink(path)


def test_is_slide_hidden(temp_pptx):
    """Test reading hidden status of slides."""
    # Create a test presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])  # Slide 1
    prs.slides.add_slide(prs.slide_layouts[1])  # Slide 2
    prs.slides.add_slide(prs.slide_layouts[0])  # Slide 3

    # Hide slide 2 - set show attribute on the slide element, not sldId
    prs.slides[1].element.set("show", "0")

    # Save to temp file
    prs.save(temp_pptx)

    # Test with PPTXHandler
    handler = PPTXHandler(temp_pptx)

    assert handler.is_slide_hidden(1) is False, "Slide 1 should be visible"
    assert handler.is_slide_hidden(2) is True, "Slide 2 should be hidden"
    assert handler.is_slide_hidden(3) is False, "Slide 3 should be visible"


def test_set_slide_hidden(temp_pptx):
    """Test setting slide visibility."""
    # Create a test presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])  # Slide 1
    prs.slides.add_slide(prs.slide_layouts[1])  # Slide 2
    prs.slides.add_slide(prs.slide_layouts[0])  # Slide 3

    # Save to temp file
    prs.save(temp_pptx)

    # Test with PPTXHandler
    handler = PPTXHandler(temp_pptx)

    # All slides should be visible initially
    assert handler.is_slide_hidden(1) is False
    assert handler.is_slide_hidden(2) is False
    assert handler.is_slide_hidden(3) is False

    # Hide slide 2
    handler.set_slide_hidden(2, True)
    assert handler.is_slide_hidden(2) is True, "Slide 2 should now be hidden"

    # Show slide 2 again
    handler.set_slide_hidden(2, False)
    assert handler.is_slide_hidden(2) is False, "Slide 2 should now be visible"

    # Hide slide 1 and 3
    handler.set_slide_hidden(1, True)
    handler.set_slide_hidden(3, True)
    assert handler.is_slide_hidden(1) is True
    assert handler.is_slide_hidden(3) is True

    # Save and reload to verify persistence
    fd, temp_file2 = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    try:
        handler.save(temp_file2)

        handler2 = PPTXHandler(temp_file2)
        assert handler2.is_slide_hidden(1) is True, "Slide 1 should remain hidden after save"
        assert handler2.is_slide_hidden(2) is False, "Slide 2 should remain visible after save"
        assert handler2.is_slide_hidden(3) is True, "Slide 3 should remain hidden after save"
    finally:
        if os.path.exists(temp_file2):
            os.unlink(temp_file2)


def test_get_slides_metadata(temp_pptx):
    """Test getting metadata for all slides."""
    # Create a test presentation
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3 = prs.slides.add_slide(prs.slide_layouts[0])

    # Set titles
    slide1.shapes.title.text = "Title 1"
    slide2.shapes.title.text = "Title 2"
    slide3.shapes.title.text = "Title 3"

    # Hide slide 2 - set show attribute on the slide element, not sldId
    prs.slides[1].element.set("show", "0")

    # Save to temp file
    prs.save(temp_pptx)

    handler = PPTXHandler(temp_pptx)

    # Get all slides metadata
    all_metadata = handler.get_slides_metadata(include_hidden=True)
    assert len(all_metadata) == 3, "Should have 3 slides"
    assert all_metadata[0]["title"] == "Title 1"
    assert all_metadata[0]["hidden"] is False
    assert all_metadata[1]["title"] == "Title 2"
    assert all_metadata[1]["hidden"] is True
    assert all_metadata[2]["title"] == "Title 3"
    assert all_metadata[2]["hidden"] is False

    # Get only visible slides
    visible_metadata = handler.get_slides_metadata(include_hidden=False)
    assert len(visible_metadata) == 2, "Should have 2 visible slides"
    assert visible_metadata[0]["title"] == "Title 1"
    assert visible_metadata[1]["title"] == "Title 3"


def test_get_slide_content_includes_hidden(temp_pptx):
    """Test that get_slide_content includes hidden status."""
    # Create a test presentation
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Test Slide"

    # Hide the slide - set show attribute on the slide element, not sldId
    prs.slides[0].element.set("show", "0")

    # Save to temp file
    prs.save(temp_pptx)

    handler = PPTXHandler(temp_pptx)

    content = handler.get_slide_content(1)
    assert "hidden" in content, "Content should include hidden field"
    assert content["hidden"] is True, "Slide should be marked as hidden"
    assert content["title"] == "Test Slide"


def test_get_presentation_info_includes_visibility_stats(temp_pptx):
    """Test that presentation info includes visibility statistics."""
    # Create a test presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.slides.add_slide(prs.slide_layouts[1])
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.slides.add_slide(prs.slide_layouts[1])

    # Hide slides 2 and 3 - set show attribute on the slide elements, not sldId
    prs.slides[1].element.set("show", "0")
    prs.slides[2].element.set("show", "0")

    # Save to temp file
    prs.save(temp_pptx)

    handler = PPTXHandler(temp_pptx)

    info = handler.get_presentation_info()
    assert "visible_slides" in info, "Info should include visible_slides"
    assert "hidden_slides" in info, "Info should include hidden_slides"
    assert info["slide_count"] == 4
    assert info["visible_slides"] == 2, "Should have 2 visible slides"
    assert info["hidden_slides"] == 2, "Should have 2 hidden slides"
