#!/usr/bin/env python3
"""Unit tests for replace_text tool."""

import tempfile
import traceback
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from mcp_server.tools.text_replace_tools import handle_replace_text


def create_test_pptx(path: Path) -> None:
    """Create a test PPTX file with sample content."""
    prs = Presentation()
    
    # Slide 1: Title slide with content
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    title.text = "Hello World"
    
    # Add notes to slide 1
    notes_slide = slide1.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "This is a test note. Hello world, hello everyone!"
    
    # Slide 2: Content slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = "Test Slide"
    
    # Add a text box
    left = Inches(1)
    top = Inches(2)
    width = Inches(5)
    height = Inches(1)
    textbox = slide2.shapes.add_textbox(left, top, width, height)
    textbox.text_frame.text = "Find this text in the slide"
    
    # Add notes to slide 2
    notes_slide2 = slide2.notes_slide
    text_frame2 = notes_slide2.notes_text_frame
    text_frame2.text = "Second slide notes with some text to find."
    
    prs.save(str(path))


import pytest

@pytest.fixture
def test_pptx_file():
    with tempfile.TemporaryDirectory() as tmpdir:
        test_file = Path(tmpdir) / "test.pptx"
        create_test_pptx(test_file)
        yield test_file

@pytest.mark.asyncio
async def test_replace_in_notes_literal(test_pptx_file):
    """Test literal text replacement in notes."""
    # Perform replacement - note that "test" only appears in slide 1 notes
    result = await handle_replace_text({
        "pptx_path": str(test_pptx_file),
        "target": "slide_notes",
        "pattern": "test",
        "replacement": "sample",
        "use_regex": False,
        "in_place": True,
    })
    
    assert result.get('success') is True
    
    # Verify changes
    prs = Presentation(str(test_pptx_file))
    slide1_notes = prs.slides[0].notes_slide.notes_text_frame.text
    slide2_notes = prs.slides[1].notes_slide.notes_text_frame.text
    
    assert "sample" in slide1_notes, "Replacement not found in slide 1 notes"
    assert "test" not in slide1_notes, "Original text still present in slide 1"
    # Slide 2 has "text" not "test", so it should be unchanged
    assert "text" in slide2_notes, "Slide 2 should be unchanged"

@pytest.mark.asyncio
async def test_replace_in_content_literal(test_pptx_file):
    """Test literal text replacement in slide content."""
    # Perform replacement
    result = await handle_replace_text({
        "pptx_path": str(test_pptx_file),
        "target": "slide_content",
        "pattern": "Hello",
        "replacement": "Goodbye",
        "use_regex": False,
        "in_place": True,
    })
    
    assert result.get('success') is True
    
    # Verify changes
    prs = Presentation(str(test_pptx_file))
    slide1_title = prs.slides[0].shapes.title.text
    
    assert "Goodbye" in slide1_title, "Replacement not found in slide 1 title"
    assert "Hello" not in slide1_title, "Original text still present"

@pytest.mark.asyncio
async def test_replace_with_regex(test_pptx_file):
    """Test regex replacement with capture groups."""
    # Perform regex replacement in notes
    result = await handle_replace_text({
        "pptx_path": str(test_pptx_file),
        "target": "slide_notes",
        "pattern": r"(hello)\s+(world)",
        "replacement": r"\2 \1",
        "use_regex": True,
        "regex_flags": ["IGNORECASE"],
        "in_place": True,
    })
    
    assert result.get('success') is True
    
    # Verify changes
    prs = Presentation(str(test_pptx_file))
    slide1_notes = prs.slides[0].notes_slide.notes_text_frame.text
    
    assert "world hello" in slide1_notes.lower(), "Regex replacement not applied correctly"

@pytest.mark.asyncio
async def test_dry_run(test_pptx_file):
    """Test dry run mode doesn't modify file."""
    # Get original content
    prs_before = Presentation(str(test_pptx_file))
    notes_before = prs_before.slides[0].notes_slide.notes_text_frame.text
    
    # Perform dry run
    result = await handle_replace_text({
        "pptx_path": str(test_pptx_file),
        "target": "slide_notes",
        "pattern": "test",
        "replacement": "sample",
        "use_regex": False,
        "dry_run": True,
    })
    
    assert result.get('success') is True
    
    # Verify file wasn't modified
    prs_after = Presentation(str(test_pptx_file))
    notes_after = prs_after.slides[0].notes_slide.notes_text_frame.text
    
    assert notes_before == notes_after, "File was modified in dry run mode!"
    assert result.get('replacements_count', 0) > 0, "No changes detected"
    assert result.get('dry_run') == True, "Dry run flag not set"

@pytest.mark.asyncio
async def test_max_replacements(test_pptx_file):
    """Test max_replacements parameter with regex."""
    # Perform replacement with limit using regex (case-insensitive)
    result = await handle_replace_text({
        "pptx_path": str(test_pptx_file),
        "target": "slide_notes",
        "pattern": "hello",
        "replacement": "hi",
        "use_regex": True,
        "regex_flags": ["IGNORECASE"],
        "max_replacements": 1,
        "in_place": True,
    })
    
    assert result.get('success') is True
    assert result.get('replacements_count') == 1, f"Expected 1 replacement, got {result.get('replacements_count')}"

@pytest.mark.asyncio
async def test_max_replacements_literal(test_pptx_file):
    """Test max_replacements parameter with literal text."""
    # Perform literal replacement with limit (note has "hello" twice)
    result = await handle_replace_text({
        "pptx_path": str(test_pptx_file),
        "target": "slide_notes",
        "pattern": "hello",
        "replacement": "hi",
        "use_regex": False,
        "max_replacements": 1,
        "in_place": True,
    })
    
    assert result.get('success') is True
    assert result.get('replacements_count') == 1, f"Expected 1 replacement, got {result.get('replacements_count')}"
    
    # Verify one "hello" remains
    prs = Presentation(str(test_pptx_file))
    notes = prs.slides[0].notes_slide.notes_text_frame.text
    assert "hi" in notes.lower(), "Should have one replacement"
    assert "hello" in notes.lower(), "Should still have one 'hello' remaining"

@pytest.mark.asyncio
async def test_specific_slide(test_pptx_file):
    """Test replacement on specific slide only."""
    # Perform replacement on slide 1 only - replace "note" with "document"
    result = await handle_replace_text({
        "pptx_path": str(test_pptx_file),
        "target": "slide_notes",
        "pattern": "note",
        "replacement": "document",
        "use_regex": False,
        "slide_number": 1,
        "in_place": True,
    })
    
    assert result.get('success') is True
    
    # Verify only slide 1 was modified
    prs = Presentation(str(test_pptx_file))
    slide1_notes = prs.slides[0].notes_slide.notes_text_frame.text
    slide2_notes = prs.slides[1].notes_slide.notes_text_frame.text
    
    assert "document" in slide1_notes, "Slide 1 not modified"
    assert "note" not in slide1_notes, "Original text still in slide 1"
    assert "notes" in slide2_notes, "Slide 2 should be unchanged"
    assert "document" not in slide2_notes, "Slide 2 was incorrectly modified"

@pytest.mark.asyncio
async def test_output_path(test_pptx_file):
    """Test output to different file (not in-place)."""
    with tempfile.TemporaryDirectory() as tmpdir:
        output_file = Path(tmpdir) / "output.pptx"
        
        # Get original content
        prs_before = Presentation(str(test_pptx_file))
        notes_before = prs_before.slides[0].notes_slide.notes_text_frame.text
        
        # Perform replacement to new file
        result = await handle_replace_text({
            "pptx_path": str(test_pptx_file),
            "target": "slide_notes",
            "pattern": "test",
            "replacement": "sample",
            "use_regex": False,
            "in_place": False,
            "output_path": str(output_file),
        })
        
        assert result.get('success') is True
        
        # Verify original file unchanged
        prs_original = Presentation(str(test_pptx_file))
        notes_original = prs_original.slides[0].notes_slide.notes_text_frame.text
        
        assert notes_before == notes_original, "Original file was modified!"
        
        # Verify output file was created and modified
        assert output_file.exists(), "Output file not created!"
        prs_output = Presentation(str(output_file))
        notes_output = prs_output.slides[0].notes_slide.notes_text_frame.text
        
        assert "sample" in notes_output, "Output file not modified correctly"

@pytest.mark.asyncio
async def test_error_invalid_regex(test_pptx_file):
    """Test error handling for invalid regex pattern."""
    # Test invalid regex pattern
    result = await handle_replace_text({
        "pptx_path": str(test_pptx_file),
        "target": "slide_notes",
        "pattern": r"([unclosed",  # Invalid regex
        "replacement": "test",
        "use_regex": True,
        "in_place": True,
    })
    
    assert result.get('success') is False
    assert "error" in result
    assert "Invalid regular expression" in result.get('error', '')

@pytest.mark.asyncio
async def test_error_invalid_file():
    """Test error handling for non-existent file."""
    # Test non-existent file
    result = await handle_replace_text({
        "pptx_path": "/nonexistent/file.pptx",
        "target": "slide_notes",
        "pattern": "test",
        "replacement": "sample",
        "use_regex": False,
        "in_place": True,
    })
    
    assert result.get('success') is False
    assert "error" in result

@pytest.mark.asyncio
async def test_error_invalid_slide_number(test_pptx_file):
    """Test error handling for invalid slide number."""
    # Test invalid slide number (file has only 2 slides)
    result = await handle_replace_text({
        "pptx_path": str(test_pptx_file),
        "target": "slide_notes",
        "pattern": "test",
        "replacement": "sample",
        "use_regex": False,
        "slide_number": 999,
        "in_place": True,
    })
    
    assert result.get('success') is False
    assert "error" in result

@pytest.mark.asyncio
async def test_error_invalid_target(test_pptx_file):
    """Test error handling for invalid target value."""
    # Test invalid target
    result = await handle_replace_text({
        "pptx_path": str(test_pptx_file),
        "target": "invalid_target",
        "pattern": "test",
        "replacement": "sample",
        "use_regex": False,
        "in_place": True,
    })
    
    assert result.get('success') is False
    assert "error" in result
    assert "Invalid target" in result.get('error', '')
