"""Core PPTX file operations handler."""

from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture

from ..cache import PresentationCache
from ..config import get_config
from ..utils.async_utils import run_in_thread
from ..utils.validators import validate_pptx_path, validate_slide_number


class PPTXHandler:
    """Handler for PPTX file operations with optional caching support.

    This handler provides core PPTX file operations with:
    - Lazy loading of presentations
    - Optional caching for better performance
    - Automatic cache invalidation on modifications
    """

    def __init__(
        self,
        pptx_path: str | Path,
        cache: Optional[PresentationCache] = None,
        enable_cache: Optional[bool] = None,
    ):
        """Initialize handler with PPTX file path.

        Args:
            pptx_path: Path to PPTX file
            cache: Optional PresentationCache instance
            enable_cache: Whether to enable caching (uses config default if None)
        """
        self.pptx_path = validate_pptx_path(pptx_path)
        self._presentation: Optional[Presentation] = None

        # Caching support
        config = get_config()
        self._enable_cache = (
            enable_cache if enable_cache is not None else config.performance.enable_cache
        )
        self._cache = cache if self._enable_cache else None
        self._is_modified = False

    async def get_presentation(self) -> Presentation:
        """Load presentation with optional caching."""
        if self._presentation is None:
            # Try to get from cache first
            if self._cache and self._enable_cache:
                cached = self._cache.get_presentation(self.pptx_path)
                if cached is not None:
                    self._presentation = cached
                    return self._presentation

            # Load from file (blocking call, run in thread)
            self._presentation = await run_in_thread(Presentation, str(self.pptx_path))

            # Cache it
            if self._cache and self._enable_cache:
                self._cache.cache_presentation(self.pptx_path, self._presentation)

        return self._presentation

    @property
    def presentation(self) -> Presentation:
        """Access presentation synchronously.

        Note: This will raise RuntimeError if the presentation hasn't been loaded yet
        via get_presentation(). Use get_presentation() in async contexts.
        """
        if self._presentation is None:
            # For backward compatibility, we'll load it synchronously but log a warning
            import logging

            logging.getLogger(__name__).warning(
                "Synchronous access to PPTXHandler.presentation is deprecated. "
                "Use await handler.get_presentation() instead."
            )
            self._presentation = Presentation(str(self.pptx_path))
        return self._presentation

    def reload(self):
        """Reload presentation from file.

        Clears cached presentation and marks for reload on next access.
        """
        self._presentation = None
        self._is_modified = False

        # Invalidate cache
        if self._cache and self._enable_cache:
            self._cache.invalidate(self.pptx_path)

    async def get_slide_count(self) -> int:
        """Get total number of slides."""
        pres = await self.get_presentation()
        return len(pres.slides)

    async def is_slide_hidden(self, slide_number: int) -> bool:
        """Check if a slide is hidden.

        Args:
            slide_number: Slide number (1-indexed)

        Returns:
            True if slide is hidden, False if visible
        """
        slide_count = await self.get_slide_count()
        validate_slide_number(slide_number, slide_count)
        pres = await self.get_presentation()
        slide = pres.slides[slide_number - 1]

        # Check for 'show' attribute in two locations:
        # 1. On the <p:sld> element (the slide part)
        # 2. On the <p:sldId> element in presentation.xml (standard PowerPoint)

        # Check location 1: <p:sld> element
        show_attr = slide.element.get("show")
        if show_attr is not None:
            if str(show_attr).strip() == "0":
                return True

        # Check location 2: <p:sldId> element
        try:
            # prs.slides._sldIdLst is the list of <p:sldId> elements
            sld_id_lst = pres.slides._sldIdLst
            if slide_number - 1 < len(sld_id_lst):
                sld_id_element = sld_id_lst[slide_number - 1]
                show_attr = sld_id_element.get("show")
                if show_attr is not None:
                    if str(show_attr).strip() == "0":
                        return True
        except (AttributeError, IndexError, ValueError):
            # Fallback if internals change or are inaccessible
            pass

        return False

    async def set_slide_hidden(self, slide_number: int, hidden: bool):
        """Set slide visibility.

        Args:
            slide_number: Slide number (1-indexed)
            hidden: True to hide the slide, False to show it
        """
        slide_count = await self.get_slide_count()
        validate_slide_number(slide_number, slide_count)
        pres = await self.get_presentation()
        slide = pres.slides[slide_number - 1]

        # Set on both locations for maximum compatibility
        if hidden:
            # Set on <p:sld>
            slide.element.set("show", "0")
            # Set on <p:sldId>
            try:
                sld_id_lst = pres.slides._sldIdLst
                sld_id_lst[slide_number - 1].set("show", "0")
            except (AttributeError, IndexError):
                pass
        else:
            # Remove from <p:sld>
            if "show" in slide.element.attrib:
                del slide.element.attrib["show"]
            # Remove from <p:sldId>
            try:
                sld_id_lst = pres.slides._sldIdLst
                sld_id_element = sld_id_lst[slide_number - 1]
                if "show" in sld_id_element.attrib:
                    del sld_id_element.attrib["show"]
            except (AttributeError, IndexError):
                pass

        # Mark as modified
        self._is_modified = True

    async def get_presentation_info(self) -> Dict[str, Any]:
        """Get presentation metadata."""
        slide_count = await self.get_slide_count()
        hidden_count = 0
        for i in range(1, slide_count + 1):
            if await self.is_slide_hidden(i):
                hidden_count += 1
        visible_count = slide_count - hidden_count

        pres = await self.get_presentation()
        return {
            "file_path": str(self.pptx_path),
            "file_name": self.pptx_path.name,
            "slide_count": slide_count,
            "visible_slides": visible_count,
            "hidden_slides": hidden_count,
            "slide_size": {
                "width": pres.slide_width,
                "height": pres.slide_height,
            },
        }

    async def get_slide_text(self, slide_number: int) -> str:
        """Extract all text from a slide."""
        slide_count = await self.get_slide_count()
        validate_slide_number(slide_number, slide_count)
        pres = await self.get_presentation()
        slide = pres.slides[slide_number - 1]
        return self._extract_text_from_slide(slide)

    def _extract_text_from_slide(self, slide) -> str:
        """Recursively extract text from all shapes in a slide."""
        texts = []
        for shape in slide.shapes:
            texts.extend(self._extract_text_from_shape(shape))
        return "\n".join(texts)

    def _extract_text_from_shape(self, shape: BaseShape) -> List[str]:
        """Recursively extract text from a shape."""
        texts = []
        if hasattr(shape, "text_frame") and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs)
                if text.strip():
                    texts.append(text)
        elif isinstance(shape, GroupShape):
            for sub_shape in shape.shapes:
                texts.extend(self._extract_text_from_shape(sub_shape))
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for row in shape.table.rows:
                row_texts = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_texts.append(cell_text)
                if row_texts:
                    texts.append(" | ".join(row_texts))
        return texts

    async def get_slide_content(self, slide_number: int) -> Dict[str, Any]:
        """Get comprehensive content from a slide."""
        slide_count = await self.get_slide_count()
        validate_slide_number(slide_number, slide_count)
        pres = await self.get_presentation()
        slide = pres.slides[slide_number - 1]

        # Extract title
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text

        # Extract text
        text = self._extract_text_from_slide(slide)

        # Extract shapes info
        shapes_info = []
        for idx, shape in enumerate(slide.shapes):
            shape_info = {
                "index": idx,
                "shape_id": shape.shape_id,
                "shape_type": str(shape.shape_type),
                "name": shape.name if hasattr(shape, "name") else "",
            }
            if hasattr(shape, "text_frame") and shape.text_frame:
                shape_info["text"] = "".join(
                    run.text for para in shape.text_frame.paragraphs for run in para.runs
                )
            if isinstance(shape, Picture):
                shape_info["image"] = True
                shape_info["image_path"] = (
                    shape.image.filename if hasattr(shape.image, "filename") else None
                )
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                shape_info["table"] = True
                shape_info["rows"] = len(shape.table.rows)
                shape_info["columns"] = len(shape.table.columns)
            shapes_info.append(shape_info)

        return {
            "slide_number": slide_number,
            "title": title,
            "text": text,
            "shapes": shapes_info,
            "hidden": await self.is_slide_hidden(slide_number),
        }

    async def get_slide_images(self, slide_number: int) -> List[Dict[str, Any]]:
        """Get image information from a slide."""
        slide_count = await self.get_slide_count()
        validate_slide_number(slide_number, slide_count)
        pres = await self.get_presentation()
        slide = pres.slides[slide_number - 1]

        images = []
        for idx, shape in enumerate(slide.shapes):
            if isinstance(shape, Picture):
                image_info = {
                    "index": idx,
                    "shape_id": shape.shape_id,
                    "name": shape.name if hasattr(shape, "name") else "",
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                }
                if hasattr(shape.image, "filename"):
                    image_info["filename"] = shape.image.filename
                images.append(image_info)
            elif isinstance(shape, GroupShape):
                # Check nested shapes in groups
                for sub_shape in shape.shapes:
                    if isinstance(sub_shape, Picture):
                        image_info = {
                            "index": idx,
                            "shape_id": sub_shape.shape_id,
                            "name": sub_shape.name if hasattr(sub_shape, "name") else "",
                            "left": sub_shape.left,
                            "top": sub_shape.top,
                            "width": sub_shape.width,
                            "height": sub_shape.height,
                        }
                        if hasattr(sub_shape.image, "filename"):
                            image_info["filename"] = sub_shape.image.filename
                        images.append(image_info)
        return images

    async def get_notes(self, slide_number: Optional[int] = None) -> Dict[str, Any]:
        """Get notes for slide(s)."""
        slide_count = await self.get_slide_count()
        pres = await self.get_presentation()

        if slide_number is not None:
            validate_slide_number(slide_number, slide_count)
            slide = pres.slides[slide_number - 1]
            notes_text = ""
            if slide.has_notes_slide:
                try:
                    notes_text = slide.notes_slide.notes_text_frame.text
                except Exception:
                    pass
            return {
                "slide": slide_number,
                "notes": notes_text,
            }
        else:
            # Return all notes
            slides = []
            for idx, slide in enumerate(pres.slides, start=1):
                notes_text = ""
                if slide.has_notes_slide:
                    try:
                        notes_text = slide.notes_slide.notes_text_frame.text
                    except Exception:
                        pass
                slides.append({"slide": idx, "notes": notes_text})
            return {"slides": slides}

    async def get_slides_metadata(self, include_hidden: bool = True) -> List[Dict[str, Any]]:
        """Get metadata for all slides including their visibility status.

        Args:
            include_hidden: If False, only return visible slides

        Returns:
            List of slide metadata dictionaries
        """
        slides_metadata = []
        slide_count = await self.get_slide_count()
        pres = await self.get_presentation()

        for i in range(1, slide_count + 1):
            is_hidden = await self.is_slide_hidden(i)

            if not include_hidden and is_hidden:
                continue

            slide = pres.slides[i - 1]
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text

            slides_metadata.append(
                {
                    "slide_number": i,
                    "title": title,
                    "hidden": is_hidden,
                    "slide_id": slide.slide_id,
                }
            )

        return slides_metadata

    async def save(self, output_path: Optional[str | Path] = None):
        """Save presentation to file.

        Args:
            output_path: Optional output path. If None, saves to original path.
        """
        save_path = Path(output_path) if output_path else self.pptx_path
        pres = await self.get_presentation()
        await run_in_thread(pres.save, str(save_path))

        # Invalidate cache after save
        if self._cache and self._enable_cache:
            self._cache.invalidate(save_path)
            # If saving to original path or same path, also invalidate original reference
            if output_path is None or Path(output_path) == self.pptx_path:
                self._cache.invalidate(self.pptx_path)

        self._is_modified = False
