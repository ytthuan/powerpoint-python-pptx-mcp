"""Tools for editing PPTX content."""

import os
import tempfile
from pathlib import Path
from typing import Any, Dict

from mcp.types import Tool
from pptx import Presentation
from pptx.util import Inches

from ..utils.validators import (
    validate_pptx_path,
    validate_slide_number,
    validate_position,
    validate_size,
    validate_image_path,
)
from ..utils.async_utils import run_in_thread


def get_edit_tools() -> list[Tool]:
    """Get all edit tools."""
    return [
        Tool(
            name="update_slide_text",
            description="Update text in a specific shape on a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "pptx_path": {
                        "type": "string",
                        "description": "Path to the PPTX file",
                    },
                    "slide_number": {
                        "type": "integer",
                        "description": "Slide number (1-indexed)",
                        "minimum": 1,
                    },
                    "shape_id": {
                        "type": "integer",
                        "description": "Shape ID to update (from read_slide_content)",
                    },
                    "new_text": {
                        "type": "string",
                        "description": "New text content",
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Output PPTX path (optional, defaults to input path with .edited suffix)",  # noqa: E501
                    },
                },
                "required": ["pptx_path", "slide_number", "shape_id", "new_text"],
            },
        ),
        Tool(
            name="replace_slide_image",
            description="Replace an existing image on a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "pptx_path": {
                        "type": "string",
                        "description": "Path to the PPTX file",
                    },
                    "slide_number": {
                        "type": "integer",
                        "description": "Slide number (1-indexed)",
                        "minimum": 1,
                    },
                    "shape_id": {
                        "type": "integer",
                        "description": "Shape ID of the image to replace",
                    },
                    "new_image_path": {
                        "type": "string",
                        "description": "Path to the new image file",
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Output PPTX path (optional)",
                    },
                },
                "required": ["pptx_path", "slide_number", "shape_id", "new_image_path"],
            },
        ),
        Tool(
            name="add_text_box",
            description="Add a new text box to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "pptx_path": {
                        "type": "string",
                        "description": "Path to the PPTX file",
                    },
                    "slide_number": {
                        "type": "integer",
                        "description": "Slide number (1-indexed)",
                        "minimum": 1,
                    },
                    "text": {
                        "type": "string",
                        "description": "Text content",
                    },
                    "position": {
                        "type": "object",
                        "properties": {
                            "x": {"type": "number", "description": "X position in inches"},
                            "y": {"type": "number", "description": "Y position in inches"},
                        },
                    },
                    "size": {
                        "type": "object",
                        "properties": {
                            "width": {"type": "number", "description": "Width in inches"},
                            "height": {"type": "number", "description": "Height in inches"},
                        },
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Output PPTX path (optional)",
                    },
                },
                "required": ["pptx_path", "slide_number", "text"],
            },
        ),
        Tool(
            name="add_image",
            description="Add a new image to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "pptx_path": {
                        "type": "string",
                        "description": "Path to the PPTX file",
                    },
                    "slide_number": {
                        "type": "integer",
                        "description": "Slide number (1-indexed)",
                        "minimum": 1,
                    },
                    "image_path": {
                        "type": "string",
                        "description": "Path to the image file",
                    },
                    "position": {
                        "type": "object",
                        "properties": {
                            "x": {"type": "number", "description": "X position in inches"},
                            "y": {"type": "number", "description": "Y position in inches"},
                        },
                    },
                    "size": {
                        "type": "object",
                        "properties": {
                            "width": {"type": "number", "description": "Width in inches"},
                            "height": {"type": "number", "description": "Height in inches"},
                        },
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Output PPTX path (optional)",
                    },
                },
                "required": ["pptx_path", "slide_number", "image_path"],
            },
        ),
        Tool(
            name="replace_slide_content",
            description="Replace all text content of a slide with new content. Clears existing text shapes and sets new content. Preserves images and non-text shapes unless clear_all is true.",  # noqa: E501
            inputSchema={
                "type": "object",
                "properties": {
                    "pptx_path": {
                        "type": "string",
                        "description": "Path to the PPTX file",
                    },
                    "slide_number": {
                        "type": "integer",
                        "description": "Slide number (1-indexed)",
                        "minimum": 1,
                    },
                    "content": {
                        "type": "object",
                        "description": "New content to set on the slide",
                        "properties": {
                            "title": {
                                "type": "string",
                                "description": "Title text (updates title shape if exists)",
                            },
                            "text_boxes": {
                                "type": "array",
                                "description": "Array of text boxes to add",
                                "items": {
                                    "type": "object",
                                    "properties": {
                                        "text": {"type": "string", "description": "Text content"},
                                        "position": {
                                            "type": "object",
                                            "properties": {
                                                "x": {
                                                    "type": "number",
                                                    "description": "X position in inches",
                                                },
                                                "y": {
                                                    "type": "number",
                                                    "description": "Y position in inches",
                                                },
                                            },
                                        },
                                        "size": {
                                            "type": "object",
                                            "properties": {
                                                "width": {
                                                    "type": "number",
                                                    "description": "Width in inches",
                                                },
                                                "height": {
                                                    "type": "number",
                                                    "description": "Height in inches",
                                                },
                                            },
                                        },
                                    },
                                    "required": ["text"],
                                },
                            },
                        },
                    },
                    "clear_all": {
                        "type": "boolean",
                        "description": "If true, removes all shapes (including images). If false, only removes text shapes (default: false)",  # noqa: E501
                        "default": False,
                    },
                    "preserve_layout": {
                        "type": "boolean",
                        "description": "If true, attempts to preserve existing shape positions when updating text (default: true)",  # noqa: E501
                        "default": True,
                    },
                    "in_place": {
                        "type": "boolean",
                        "description": "Update file in-place. If false, creates a new file (default: false)",  # noqa: E501
                        "default": False,
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Output PPTX path (only used when in_place is false)",
                    },
                },
                "required": ["pptx_path", "slide_number", "content"],
            },
        ),
        Tool(
            name="update_slide_content",
            description="Update specific parts of slide content without replacing everything. Can update title, specific shapes by ID, or add new content.",  # noqa: E501
            inputSchema={
                "type": "object",
                "properties": {
                    "pptx_path": {
                        "type": "string",
                        "description": "Path to the PPTX file",
                    },
                    "slide_number": {
                        "type": "integer",
                        "description": "Slide number (1-indexed)",
                        "minimum": 1,
                    },
                    "updates": {
                        "type": "object",
                        "description": "Content updates to apply",
                        "properties": {
                            "title": {
                                "type": "string",
                                "description": "New title text (updates title shape if exists, or creates one if not)",  # noqa: E501
                            },
                            "shape_updates": {
                                "type": "array",
                                "description": "Updates for specific shapes by shape_id",
                                "items": {
                                    "type": "object",
                                    "properties": {
                                        "shape_id": {
                                            "type": "integer",
                                            "description": "Shape ID to update (from read_slide_content)",  # noqa: E501
                                        },
                                        "text": {
                                            "type": "string",
                                            "description": "New text content for the shape",
                                        },
                                    },
                                    "required": ["shape_id", "text"],
                                },
                            },
                            "add_text_boxes": {
                                "type": "array",
                                "description": "New text boxes to add to the slide",
                                "items": {
                                    "type": "object",
                                    "properties": {
                                        "text": {"type": "string", "description": "Text content"},
                                        "position": {
                                            "type": "object",
                                            "properties": {
                                                "x": {
                                                    "type": "number",
                                                    "description": "X position in inches",
                                                },
                                                "y": {
                                                    "type": "number",
                                                    "description": "Y position in inches",
                                                },
                                            },
                                        },
                                        "size": {
                                            "type": "object",
                                            "properties": {
                                                "width": {
                                                    "type": "number",
                                                    "description": "Width in inches",
                                                },
                                                "height": {
                                                    "type": "number",
                                                    "description": "Height in inches",
                                                },
                                            },
                                        },
                                    },
                                    "required": ["text"],
                                },
                            },
                        },
                    },
                    "in_place": {
                        "type": "boolean",
                        "description": "Update file in-place. If false, creates a new file (default: false)",  # noqa: E501
                        "default": False,
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Output PPTX path (only used when in_place is false)",
                    },
                },
                "required": ["pptx_path", "slide_number", "updates"],
            },
        ),
    ]


async def handle_update_slide_text(arguments: Dict[str, Any]) -> Dict[str, Any]:
    """Handle update_slide_text tool call."""
    pptx_path = validate_pptx_path(arguments["pptx_path"])
    slide_number = arguments["slide_number"]
    shape_id = arguments["shape_id"]
    new_text = arguments["new_text"]
    output_path = arguments.get("output_path")

    pres = await run_in_thread(Presentation, str(pptx_path))
    validate_slide_number(slide_number, len(pres.slides))

    slide = pres.slides[slide_number - 1]

    # Find shape by ID
    shape = None
    for s in slide.shapes:
        if s.shape_id == shape_id:
            shape = s
            break

    if not shape:
        raise ValueError(f"Shape with ID {shape_id} not found on slide {slide_number}")

    if not hasattr(shape, "text_frame") or not shape.text_frame:
        raise ValueError(f"Shape {shape_id} does not have a text frame")

    # Update text
    shape.text_frame.text = new_text

    # Save
    if output_path:
        output_path = Path(output_path)
    else:
        output_path = pptx_path.with_name(pptx_path.stem + ".edited.pptx")

    await run_in_thread(pres.save, str(output_path))

    return {
        "success": True,
        "output_path": str(output_path),
        "slide_number": slide_number,
        "shape_id": shape_id,
    }


async def handle_replace_slide_image(arguments: Dict[str, Any]) -> Dict[str, Any]:
    """Handle replace_slide_image tool call."""
    pptx_path = validate_pptx_path(arguments["pptx_path"])
    slide_number = arguments["slide_number"]
    shape_id = arguments["shape_id"]
    new_image_path = validate_image_path(arguments["new_image_path"])
    output_path = arguments.get("output_path")

    pres = await run_in_thread(Presentation, str(pptx_path))
    validate_slide_number(slide_number, len(pres.slides))

    slide = pres.slides[slide_number - 1]

    # Find shape by ID
    from pptx.shapes.picture import Picture

    shape = None
    for s in slide.shapes:
        if s.shape_id == shape_id:
            if isinstance(s, Picture):
                shape = s
                break

    if not shape:
        raise ValueError(f"Image shape with ID {shape_id} not found on slide {slide_number}")

    # Get position and size
    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height

    # Remove old shape
    sp = shape._element
    sp.getparent().remove(sp)

    # Add new image
    slide.shapes.add_picture(str(new_image_path), left, top, width, height)

    # Save
    if output_path:
        output_path = Path(output_path)
    else:
        output_path = pptx_path.with_name(pptx_path.stem + ".edited.pptx")

    await run_in_thread(pres.save, str(output_path))

    return {
        "success": True,
        "output_path": str(output_path),
        "slide_number": slide_number,
    }


async def handle_add_text_box(arguments: Dict[str, Any]) -> Dict[str, Any]:
    """Handle add_text_box tool call."""
    pptx_path = validate_pptx_path(arguments["pptx_path"])
    slide_number = arguments["slide_number"]
    text = arguments["text"]
    position = validate_position(arguments.get("position"))
    size = validate_size(arguments.get("size"))
    output_path = arguments.get("output_path")

    pres = await run_in_thread(Presentation, str(pptx_path))
    validate_slide_number(slide_number, len(pres.slides))

    slide = pres.slides[slide_number - 1]

    # Add text box
    left = Inches(position["x"])
    top = Inches(position["y"])
    width = Inches(size["width"])
    height = Inches(size["height"])

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = text

    # Save
    if output_path:
        output_path = Path(output_path)
    else:
        output_path = pptx_path.with_name(pptx_path.stem + ".edited.pptx")

    await run_in_thread(pres.save, str(output_path))

    return {
        "success": True,
        "output_path": str(output_path),
        "slide_number": slide_number,
        "shape_id": text_box.shape_id,
    }


async def handle_add_image(arguments: Dict[str, Any]) -> Dict[str, Any]:
    """Handle add_image tool call."""
    pptx_path = validate_pptx_path(arguments["pptx_path"])
    slide_number = arguments["slide_number"]
    image_path = validate_image_path(arguments["image_path"])
    position = validate_position(arguments.get("position"))
    size = validate_size(arguments.get("size"))
    output_path = arguments.get("output_path")

    pres = await run_in_thread(Presentation, str(pptx_path))
    validate_slide_number(slide_number, len(pres.slides))

    slide = pres.slides[slide_number - 1]

    # Add image
    left = Inches(position["x"])
    top = Inches(position["y"])
    width = Inches(size["width"])
    height = Inches(size["height"])

    picture = slide.shapes.add_picture(str(image_path), left, top, width, height)

    # Save
    if output_path:
        output_path = Path(output_path)
    else:
        output_path = pptx_path.with_name(pptx_path.stem + ".edited.pptx")

    await run_in_thread(pres.save, str(output_path))

    return {
        "success": True,
        "output_path": str(output_path),
        "slide_number": slide_number,
        "shape_id": picture.shape_id,
    }


async def handle_replace_slide_content(arguments: Dict[str, Any]) -> Dict[str, Any]:
    """Handle replace_slide_content tool call."""
    pptx_path = validate_pptx_path(arguments["pptx_path"])
    slide_number = arguments["slide_number"]
    content = arguments["content"]
    clear_all = arguments.get("clear_all", False)
    preserve_layout = arguments.get("preserve_layout", True)
    in_place = arguments.get("in_place", False)
    output_path = arguments.get("output_path")

    pres = await run_in_thread(Presentation, str(pptx_path))
    validate_slide_number(slide_number, len(pres.slides))

    slide = pres.slides[slide_number - 1]

    # Store positions of existing shapes if preserving layout
    shape_positions = {}
    if preserve_layout and not clear_all:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                shape_positions[shape.shape_id] = {
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                }

    # Remove shapes
    shapes_to_remove = []
    for shape in slide.shapes:
        if clear_all:
            # Remove all shapes
            shapes_to_remove.append(shape)
        else:
            # Only remove text shapes (preserve images and other non-text shapes)
            if hasattr(shape, "text_frame") and shape.text_frame:
                # Don't remove title shape yet if we're going to update it
                if hasattr(slide.shapes, "title") and shape == slide.shapes.title:
                    if "title" in content:
                        continue  # Keep it to update
                shapes_to_remove.append(shape)

    # Remove shapes (in reverse order to avoid index issues)
    for shape in reversed(shapes_to_remove):
        sp = shape._element
        sp.getparent().remove(sp)

    # Add new content
    updated_shapes = []

    # Update or add title
    if "title" in content and content["title"]:
        if hasattr(slide.shapes, "title") and slide.shapes.title:
            slide.shapes.title.text = content["title"]
            updated_shapes.append({"type": "title", "shape_id": slide.shapes.title.shape_id})
        else:
            # Add title text box
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_box.text_frame.text = content["title"]
            updated_shapes.append({"type": "title", "shape_id": title_box.shape_id})

    # Add text boxes
    if "text_boxes" in content and content["text_boxes"]:
        for text_box_spec in content["text_boxes"]:
            text = text_box_spec["text"]
            position = validate_position(text_box_spec.get("position"))
            size = validate_size(text_box_spec.get("size"))

            left = Inches(position["x"])
            top = Inches(position["y"])
            width = Inches(size["width"])
            height = Inches(size["height"])

            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_box.text_frame.text = text
            updated_shapes.append({"type": "text_box", "shape_id": text_box.shape_id})

    # Save
    if in_place:
        # Save to temp file first, then replace
        tmp_dir = str(pptx_path.parent)
        fd, tmp_path = tempfile.mkstemp(
            prefix=pptx_path.stem + ".", suffix=".tmp.pptx", dir=tmp_dir
        )
        os.close(fd)
        tmp_file = Path(tmp_path)
        try:
            await run_in_thread(pres.save, str(tmp_file))
            os.replace(str(tmp_file), str(pptx_path))
            final_path = str(pptx_path)
        finally:
            if tmp_file.exists():
                try:
                    tmp_file.unlink()
                except Exception:
                    pass
    else:
        if output_path:
            output_path = Path(output_path)
        else:
            output_path = pptx_path.with_name(pptx_path.stem + ".replaced.pptx")
        await run_in_thread(pres.save, str(output_path))
        final_path = str(output_path)

    return {
        "success": True,
        "output_path": final_path,
        "slide_number": slide_number,
        "in_place": in_place,
        "shapes_updated": len(updated_shapes),
        "updated_shapes": updated_shapes,
    }


async def handle_update_slide_content(arguments: Dict[str, Any]) -> Dict[str, Any]:
    """Handle update_slide_content tool call."""
    pptx_path = validate_pptx_path(arguments["pptx_path"])
    slide_number = arguments["slide_number"]
    updates = arguments["updates"]
    in_place = arguments.get("in_place", False)
    output_path = arguments.get("output_path")

    pres = await run_in_thread(Presentation, str(pptx_path))
    validate_slide_number(slide_number, len(pres.slides))

    slide = pres.slides[slide_number - 1]

    updated_shapes = []
    errors = []

    # Update title if specified
    if "title" in updates and updates["title"]:
        if hasattr(slide.shapes, "title") and slide.shapes.title:
            slide.shapes.title.text = updates["title"]
            updated_shapes.append({"type": "title", "shape_id": slide.shapes.title.shape_id})
        else:
            # Add title text box if it doesn't exist
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_box.text_frame.text = updates["title"]
            updated_shapes.append({"type": "title", "shape_id": title_box.shape_id})

    # Update specific shapes by ID
    if "shape_updates" in updates and updates["shape_updates"]:
        for shape_update in updates["shape_updates"]:
            shape_id = shape_update["shape_id"]
            new_text = shape_update["text"]

            # Find shape by ID
            shape = None
            for s in slide.shapes:
                if s.shape_id == shape_id:
                    shape = s
                    break

            if not shape:
                errors.append(f"Shape with ID {shape_id} not found")
                continue

            if not hasattr(shape, "text_frame") or not shape.text_frame:
                errors.append(f"Shape {shape_id} does not have a text frame")
                continue

            # Update text
            shape.text_frame.text = new_text
            updated_shapes.append({"type": "shape_update", "shape_id": shape_id})

    # Add new text boxes
    if "add_text_boxes" in updates and updates["add_text_boxes"]:
        for text_box_spec in updates["add_text_boxes"]:
            text = text_box_spec["text"]
            position = validate_position(text_box_spec.get("position"))
            size = validate_size(text_box_spec.get("size"))

            left = Inches(position["x"])
            top = Inches(position["y"])
            width = Inches(size["width"])
            height = Inches(size["height"])

            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_box.text_frame.text = text
            updated_shapes.append({"type": "text_box", "shape_id": text_box.shape_id})

    # Save
    if in_place:
        # Save to temp file first, then replace
        tmp_dir = str(pptx_path.parent)
        fd, tmp_path = tempfile.mkstemp(
            prefix=pptx_path.stem + ".", suffix=".tmp.pptx", dir=tmp_dir
        )
        os.close(fd)
        tmp_file = Path(tmp_path)
        try:
            await run_in_thread(pres.save, str(tmp_file))
            os.replace(str(tmp_file), str(pptx_path))
            final_path = str(pptx_path)
        finally:
            if tmp_file.exists():
                try:
                    tmp_file.unlink()
                except Exception:
                    pass
    else:
        if output_path:
            output_path = Path(output_path)
        else:
            output_path = pptx_path.with_name(pptx_path.stem + ".updated.pptx")
        await run_in_thread(pres.save, str(output_path))
        final_path = str(output_path)

    result = {
        "success": True,
        "output_path": final_path,
        "slide_number": slide_number,
        "in_place": in_place,
        "shapes_updated": len(updated_shapes),
        "updated_shapes": updated_shapes,
    }

    if errors:
        result["errors"] = errors
        result["success"] = len(errors) < len(updates.get("shape_updates", []))

    return result
