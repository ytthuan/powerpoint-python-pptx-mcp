"""Tools for text replacement in PPTX files (slide content and notes)."""

import logging
import os
import re
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Tuple, Optional

from mcp.types import Tool
from pptx import Presentation

from ..core.pptx_handler import PPTXHandler
from ..core.safe_editor import update_notes_safe, update_notes_safe_in_place
from ..utils.validators import validate_pptx_path, validate_slide_number
from ..utils.async_utils import run_in_thread

logger = logging.getLogger(__name__)


def get_text_replace_tools() -> list[Tool]:
    """Get all text replacement tools."""
    return [
        Tool(
            name="replace_text",
            description="Replace text in slide content or speaker notes with optional regex support. Supports dry-run mode to preview changes without modifying files.",  # noqa: E501
            inputSchema={
                "type": "object",
                "properties": {
                    "pptx_path": {
                        "type": "string",
                        "description": "Path to the PPTX file",
                    },
                    "target": {
                        "type": "string",
                        "enum": ["slide_notes", "slide_content"],
                        "description": "Target for replacement: 'slide_notes' (speaker notes) or 'slide_content' (text in slide shapes)",  # noqa: E501
                    },
                    "pattern": {
                        "type": "string",
                        "description": "Text pattern to find. If use_regex is true, this is a regex pattern; otherwise, literal text.",  # noqa: E501
                    },
                    "replacement": {
                        "type": "string",
                        "description": "Replacement text. For regex mode, can use capture groups like \\1, \\2, etc.",  # noqa: E501
                    },
                    "use_regex": {
                        "type": "boolean",
                        "description": "Enable regex pattern matching (default: false for literal text matching)",  # noqa: E501
                        "default": False,
                    },
                    "regex_flags": {
                        "type": "array",
                        "description": "Regex flags: 'IGNORECASE', 'MULTILINE', 'DOTALL' (only used if use_regex is true)",  # noqa: E501
                        "items": {
                            "type": "string",
                            "enum": ["IGNORECASE", "MULTILINE", "DOTALL"],
                        },
                    },
                    "slide_number": {
                        "type": "integer",
                        "description": "Specific slide number (1-indexed). If omitted, applies to all slides.",  # noqa: E501
                        "minimum": 1,
                    },
                    "shape_id": {
                        "type": "integer",
                        "description": "Specific shape ID (only for slide_content target). If omitted, applies to all shapes.",  # noqa: E501
                    },
                    "max_replacements": {
                        "type": "integer",
                        "description": "Maximum number of replacements to make (0 or omitted = replace all)",  # noqa: E501
                        "minimum": 0,
                    },
                    "dry_run": {
                        "type": "boolean",
                        "description": "Preview changes without modifying the file (default: false)",  # noqa: E501
                        "default": False,
                    },
                    "in_place": {
                        "type": "boolean",
                        "description": "Update file in-place. If false, creates a new file (default: true)",  # noqa: E501
                        "default": True,
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Output path for the modified file (only used when in_place is false)",  # noqa: E501
                    },
                },
                "required": ["pptx_path", "target", "pattern", "replacement"],
            },
        ),
    ]


def _compile_regex(pattern: str, regex_flags: Optional[List[str]] = None) -> re.Pattern:
    """Compile regex pattern with specified flags."""
    flags = 0
    if regex_flags:
        for flag_name in regex_flags:
            if flag_name == "IGNORECASE":
                flags |= re.IGNORECASE
            elif flag_name == "MULTILINE":
                flags |= re.MULTILINE
            elif flag_name == "DOTALL":
                flags |= re.DOTALL

    try:
        return re.compile(pattern, flags)
    except re.error as exc:
        # Provide a clear, user-facing error message for invalid regex patterns
        raise ValueError(f"Invalid regular expression pattern {pattern!r}: {exc}") from exc


def _replace_in_text(
    text: str,
    pattern: str,
    replacement: str,
    use_regex: bool,
    regex_flags: Optional[List[str]] = None,
    max_replacements: int = 0,
) -> Tuple[str, int]:
    """Replace text with pattern, return (new_text, replacement_count)."""
    if not text:
        return text, 0

    count = 0
    if use_regex:
        regex = _compile_regex(pattern, regex_flags)
        if max_replacements > 0:
            new_text, count = regex.subn(replacement, text, count=max_replacements)
        else:
            new_text, count = regex.subn(replacement, text)
    else:
        # Literal text replacement
        if max_replacements > 0:
            # Count occurrences before replacement
            count = min(text.count(pattern), max_replacements)
            new_text = text.replace(pattern, replacement, max_replacements)
        else:
            count = text.count(pattern)
            new_text = text.replace(pattern, replacement)

    return new_text, count


async def _replace_in_notes(
    pptx_path: Path,
    pattern: str,
    replacement: str,
    use_regex: bool,
    regex_flags: Optional[List[str]],
    slide_number: Optional[int],
    max_replacements: int,
    dry_run: bool,
) -> Dict[str, Any]:
    """Replace text in speaker notes."""
    handler = PPTXHandler(pptx_path)
    max_slides = await handler.get_slide_count()

    # Determine which slides to process
    if slide_number:
        validate_slide_number(slide_number, max_slides)
        slide_numbers = [slide_number]
    else:
        slide_numbers = list(range(1, max_slides + 1))

    # Track changes
    changes: List[Dict[str, Any]] = []
    updates: List[Tuple[int, str]] = []
    total_replacements = 0
    slides_changed = 0
    remaining_replacements = max_replacements if max_replacements > 0 else float("inf")

    pres = await handler.get_presentation()
    for slide_num in slide_numbers:
        if remaining_replacements <= 0:
            break

        # Get current notes
        slide = pres.slides[slide_num - 1]
        notes_text = ""
        if slide.has_notes_slide:
            try:
                notes_text = slide.notes_slide.notes_text_frame.text
            except Exception:
                continue

        if not notes_text:
            continue

        # Perform replacement
        current_max = int(remaining_replacements) if remaining_replacements != float("inf") else 0
        new_notes, count = _replace_in_text(
            notes_text, pattern, replacement, use_regex, regex_flags, current_max
        )

        if count > 0:
            changes.append(
                {
                    "slide_number": slide_num,
                    "replacements": count,
                    "original_preview": (
                        notes_text[:100] + "..." if len(notes_text) > 100 else notes_text
                    ),
                    "modified_preview": (
                        new_notes[:100] + "..." if len(new_notes) > 100 else new_notes
                    ),
                }
            )
            updates.append((slide_num, new_notes))
            total_replacements += count
            slides_changed += 1

            if max_replacements > 0:
                remaining_replacements -= count

    return {
        "updates": updates,
        "changes": changes,
        "total_replacements": total_replacements,
        "slides_changed": slides_changed,
        "slides_scanned": len(slide_numbers),
    }


def _replace_in_content(
    pptx_path: Path,
    pattern: str,
    replacement: str,
    use_regex: bool,
    regex_flags: Optional[List[str]],
    slide_number: Optional[int],
    shape_id: Optional[int],
    max_replacements: int,
    dry_run: bool,
) -> Tuple[Presentation, Dict[str, Any]]:
    """Replace text in slide content shapes."""
    pres = Presentation(str(pptx_path))
    max_slides = len(pres.slides)

    # Determine which slides to process
    if slide_number:
        validate_slide_number(slide_number, max_slides)
        slide_numbers = [slide_number]
    else:
        slide_numbers = list(range(1, max_slides + 1))

    # Track changes
    changes: List[Dict[str, Any]] = []
    total_replacements = 0
    slides_changed = 0
    remaining_replacements = max_replacements if max_replacements > 0 else float("inf")

    for slide_num in slide_numbers:
        if remaining_replacements <= 0:
            break

        slide = pres.slides[slide_num - 1]
        slide_changed = False

        for shape in slide.shapes:
            if remaining_replacements <= 0:
                break

            # Filter by shape_id if specified
            if shape_id is not None and shape.shape_id != shape_id:
                continue

            # Only process shapes with text frames
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                continue

            # Process each paragraph in the text frame
            for paragraph in shape.text_frame.paragraphs:
                if remaining_replacements <= 0:
                    break

                # Get the full paragraph text
                para_text = "".join(run.text for run in paragraph.runs)
                if not para_text:
                    continue

                # Perform replacement
                current_max = (
                    int(remaining_replacements) if remaining_replacements != float("inf") else 0
                )
                new_para_text, count = _replace_in_text(
                    para_text, pattern, replacement, use_regex, regex_flags, current_max
                )

                if count > 0:
                    if not slide_changed:
                        slides_changed += 1
                        slide_changed = True

                    # Record change
                    changes.append(
                        {
                            "slide_number": slide_num,
                            "shape_id": shape.shape_id,
                            "replacements": count,
                            "original_text": para_text,
                            "modified_text": new_para_text,
                        }
                    )

                    # Update the paragraph text by clearing and setting the first run
                    # Note: This approach preserves the formatting of the first run but
                    # loses formatting from other runs if the paragraph had multiple runs
                    # with different formatting (e.g., bold, italic, colors).
                    # This is a known limitation for simplicity and performance.
                    if not dry_run:
                        # Clear all runs and set text in first run
                        for run in paragraph.runs[1:]:
                            run.text = ""
                        if paragraph.runs:
                            paragraph.runs[0].text = new_para_text

                    total_replacements += count

                    if max_replacements > 0:
                        remaining_replacements -= count

    return pres, {
        "changes": changes,
        "total_replacements": total_replacements,
        "slides_changed": slides_changed,
        "slides_scanned": len(slide_numbers),
    }


async def handle_replace_text(arguments: Dict[str, Any]) -> Dict[str, Any]:
    """Handle replace_text tool call."""
    try:
        # Parse arguments
        pptx_path = validate_pptx_path(arguments["pptx_path"])
        target = arguments["target"]
        pattern = arguments["pattern"]
        replacement = arguments["replacement"]
        use_regex = arguments.get("use_regex", False)
        regex_flags = arguments.get("regex_flags")
        slide_number = arguments.get("slide_number")
        shape_id = arguments.get("shape_id")
        max_replacements = arguments.get("max_replacements", 0)
        dry_run = arguments.get("dry_run", False)
        in_place = arguments.get("in_place", True)
        output_path = arguments.get("output_path")

        # Validate shape_id only for slide_content
        if shape_id is not None and target != "slide_content":
            return {
                "success": False,
                "error": "shape_id parameter is only valid for target='slide_content'",
            }

        # Process based on target
        if target == "slide_notes":
            result = await _replace_in_notes(
                pptx_path,
                pattern,
                replacement,
                use_regex,
                regex_flags,
                slide_number,
                max_replacements,
                dry_run,
            )

            if dry_run or result["total_replacements"] == 0:
                return {
                    "success": True,
                    "dry_run": dry_run,
                    "target": target,
                    "pptx_path": str(pptx_path),
                    "slides_scanned": result["slides_scanned"],
                    "slides_changed": result["slides_changed"],
                    "replacements_count": result["total_replacements"],
                    "affected_slides": [c["slide_number"] for c in result["changes"]],
                    "changes": result["changes"],
                }

            # Apply changes
            if in_place:
                await run_in_thread(update_notes_safe_in_place, pptx_path, result["updates"])
                final_path = str(pptx_path)
            else:
                if output_path:
                    output_path = Path(output_path)
                else:
                    output_path = pptx_path.with_name(pptx_path.stem + ".replaced.pptx")
                await run_in_thread(update_notes_safe, pptx_path, result["updates"], output_path)
                final_path = str(output_path)

            return {
                "success": True,
                "target": target,
                "pptx_path": final_path,
                "in_place": in_place,
                "slides_scanned": result["slides_scanned"],
                "slides_changed": result["slides_changed"],
                "replacements_count": result["total_replacements"],
                "affected_slides": [c["slide_number"] for c in result["changes"]],
            }

        elif target == "slide_content":
            # Run the blocking replacement logic in a thread
            pres, result = await run_in_thread(
                _replace_in_content,
                pptx_path,
                pattern,
                replacement,
                use_regex,
                regex_flags,
                slide_number,
                shape_id,
                max_replacements,
                dry_run,
            )

            if dry_run or result["total_replacements"] == 0:
                return {
                    "success": True,
                    "dry_run": dry_run,
                    "target": target,
                    "pptx_path": str(pptx_path),
                    "slides_scanned": result["slides_scanned"],
                    "slides_changed": result["slides_changed"],
                    "replacements_count": result["total_replacements"],
                    "affected_shapes": [
                        (c["slide_number"], c["shape_id"]) for c in result["changes"]
                    ],
                    "changes": result["changes"],
                }

            # Save changes
            if in_place:
                # Save to temp file first, then replace
                tmp_dir = str(pptx_path.parent)
                fd, tmp_path = tempfile.mkstemp(
                    prefix=pptx_path.stem + ".", suffix=".tmp.pptx", dir=tmp_dir
                )
                os.close(fd)
                tmp_file = Path(tmp_path)
                try:
                    await run_in_thread(pres.save, str(tmp_file))
                    os.replace(str(tmp_file), str(pptx_path))
                    final_path = str(pptx_path)
                finally:
                    if tmp_file.exists():
                        try:
                            tmp_file.unlink()
                        except Exception as cleanup_error:
                            logger.warning(
                                f"Failed to cleanup temporary file {tmp_file}: {cleanup_error}"
                            )
            else:
                if output_path:
                    output_path = Path(output_path)
                else:
                    output_path = pptx_path.with_name(pptx_path.stem + ".replaced.pptx")
                await run_in_thread(pres.save, str(output_path))
                final_path = str(output_path)

            return {
                "success": True,
                "target": target,
                "pptx_path": final_path,
                "in_place": in_place,
                "slides_scanned": result["slides_scanned"],
                "slides_changed": result["slides_changed"],
                "replacements_count": result["total_replacements"],
                "affected_shapes": [(c["slide_number"], c["shape_id"]) for c in result["changes"]],
            }

        else:
            return {
                "success": False,
                "error": f"Invalid target: {target}. Must be 'slide_notes' or 'slide_content'",
            }

    except ValueError as e:
        # Validation errors, regex errors, etc.
        return {
            "success": False,
            "error": str(e),
        }
    except FileNotFoundError as e:
        return {
            "success": False,
            "error": f"File not found: {e}",
        }
    except PermissionError as e:
        return {
            "success": False,
            "error": f"Permission denied: {e}",
        }
    except Exception as e:
        # Log unexpected errors for debugging
        logger.error(f"Unexpected error in replace_text: {e}", exc_info=True)
        return {
            "success": False,
            "error": f"Unexpected error: {str(e)}",
        }
