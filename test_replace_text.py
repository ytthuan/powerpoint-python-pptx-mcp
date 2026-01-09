#!/usr/bin/env python3
"""Unit tests for replace_text tool."""

import os
import tempfile
from pathlib import Path
import shutil

from pptx import Presentation
from pptx.util import Inches, Pt

from mcp_server.tools.text_replace_tools import handle_replace_text


def create_test_pptx(path: Path) -> None:
    """Create a test PPTX file with sample content."""
    prs = Presentation()
    
    # Slide 1: Title slide with content
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    title.text = "Hello World"
    
    # Add notes to slide 1
    notes_slide = slide1.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "This is a test note. Hello world, hello everyone!"
    
    # Slide 2: Content slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = "Test Slide"
    
    # Add a text box
    left = Inches(1)
    top = Inches(2)
    width = Inches(5)
    height = Inches(1)
    textbox = slide2.shapes.add_textbox(left, top, width, height)
    textbox.text_frame.text = "Find this text in the slide"
    
    # Add notes to slide 2
    notes_slide2 = slide2.notes_slide
    text_frame2 = notes_slide2.notes_text_frame
    text_frame2.text = "Second slide notes with some text to find."
    
    prs.save(str(path))


async def test_replace_in_notes_literal():
    """Test literal text replacement in notes."""
    print("\n" + "="*60)
    print("Test: Replace in notes (literal)")
    print("="*60)
    
    with tempfile.TemporaryDirectory() as tmpdir:
        test_file = Path(tmpdir) / "test.pptx"
        create_test_pptx(test_file)
        
        # Perform replacement - note that "test" only appears in slide 1 notes
        result = await handle_replace_text({
            "pptx_path": str(test_file),
            "target": "slide_notes",
            "pattern": "test",
            "replacement": "sample",
            "use_regex": False,
            "in_place": True,
        })
        
        print(f"Success: {result.get('success')}")
        print(f"Slides scanned: {result.get('slides_scanned')}")
        print(f"Slides changed: {result.get('slides_changed')}")
        print(f"Total replacements: {result.get('replacements_count')}")
        print(f"Affected slides: {result.get('affected_slides')}")
        
        # Verify changes
        prs = Presentation(str(test_file))
        slide1_notes = prs.slides[0].notes_slide.notes_text_frame.text
        slide2_notes = prs.slides[1].notes_slide.notes_text_frame.text
        
        print(f"\nSlide 1 notes (after): {slide1_notes}")
        print(f"Slide 2 notes (after): {slide2_notes}")
        
        assert "sample" in slide1_notes, "Replacement not found in slide 1 notes"
        assert "test" not in slide1_notes, "Original text still present in slide 1"
        # Slide 2 has "text" not "test", so it should be unchanged
        assert "text" in slide2_notes, "Slide 2 should be unchanged"
        
        print("\n✅ Test passed!")


async def test_replace_in_content_literal():
    """Test literal text replacement in slide content."""
    print("\n" + "="*60)
    print("Test: Replace in content (literal)")
    print("="*60)
    
    with tempfile.TemporaryDirectory() as tmpdir:
        test_file = Path(tmpdir) / "test.pptx"
        create_test_pptx(test_file)
        
        # Perform replacement
        result = await handle_replace_text({
            "pptx_path": str(test_file),
            "target": "slide_content",
            "pattern": "Hello",
            "replacement": "Goodbye",
            "use_regex": False,
            "in_place": True,
        })
        
        print(f"Success: {result.get('success')}")
        print(f"Slides scanned: {result.get('slides_scanned')}")
        print(f"Slides changed: {result.get('slides_changed')}")
        print(f"Total replacements: {result.get('replacements_count')}")
        print(f"Affected shapes: {result.get('affected_shapes')}")
        
        # Verify changes
        prs = Presentation(str(test_file))
        slide1_title = prs.slides[0].shapes.title.text
        
        print(f"\nSlide 1 title (after): {slide1_title}")
        
        assert "Goodbye" in slide1_title, "Replacement not found in slide 1 title"
        assert "Hello" not in slide1_title, "Original text still present"
        
        print("\n✅ Test passed!")


async def test_replace_with_regex():
    """Test regex replacement with capture groups."""
    print("\n" + "="*60)
    print("Test: Replace with regex")
    print("="*60)
    
    with tempfile.TemporaryDirectory() as tmpdir:
        test_file = Path(tmpdir) / "test.pptx"
        create_test_pptx(test_file)
        
        # Perform regex replacement in notes
        result = await handle_replace_text({
            "pptx_path": str(test_file),
            "target": "slide_notes",
            "pattern": r"(hello)\s+(world)",
            "replacement": r"\2 \1",
            "use_regex": True,
            "regex_flags": ["IGNORECASE"],
            "in_place": True,
        })
        
        print(f"Success: {result.get('success')}")
        print(f"Total replacements: {result.get('replacements_count')}")
        
        # Verify changes
        prs = Presentation(str(test_file))
        slide1_notes = prs.slides[0].notes_slide.notes_text_frame.text
        
        print(f"\nSlide 1 notes (after): {slide1_notes}")
        
        assert "world hello" in slide1_notes.lower(), "Regex replacement not applied correctly"
        
        print("\n✅ Test passed!")


async def test_dry_run():
    """Test dry run mode doesn't modify file."""
    print("\n" + "="*60)
    print("Test: Dry run mode")
    print("="*60)
    
    with tempfile.TemporaryDirectory() as tmpdir:
        test_file = Path(tmpdir) / "test.pptx"
        create_test_pptx(test_file)
        
        # Get original content
        prs_before = Presentation(str(test_file))
        notes_before = prs_before.slides[0].notes_slide.notes_text_frame.text
        
        # Perform dry run
        result = await handle_replace_text({
            "pptx_path": str(test_file),
            "target": "slide_notes",
            "pattern": "test",
            "replacement": "sample",
            "use_regex": False,
            "dry_run": True,
        })
        
        print(f"Success: {result.get('success')}")
        print(f"Dry run: {result.get('dry_run')}")
        print(f"Changes found: {result.get('replacements_count')}")
        print(f"Changes detail: {result.get('changes')}")
        
        # Verify file wasn't modified
        prs_after = Presentation(str(test_file))
        notes_after = prs_after.slides[0].notes_slide.notes_text_frame.text
        
        assert notes_before == notes_after, "File was modified in dry run mode!"
        assert result.get('replacements_count', 0) > 0, "No changes detected"
        assert result.get('dry_run') == True, "Dry run flag not set"
        
        print("\n✅ Test passed!")


async def test_max_replacements():
    """Test max_replacements parameter."""
    print("\n" + "="*60)
    print("Test: Max replacements limit")
    print("="*60)
    
    with tempfile.TemporaryDirectory() as tmpdir:
        test_file = Path(tmpdir) / "test.pptx"
        create_test_pptx(test_file)
        
        # Perform replacement with limit
        result = await handle_replace_text({
            "pptx_path": str(test_file),
            "target": "slide_notes",
            "pattern": "hello",
            "replacement": "hi",
            "use_regex": False,
            "regex_flags": ["IGNORECASE"],
            "max_replacements": 1,
            "in_place": True,
        })
        
        print(f"Success: {result.get('success')}")
        print(f"Total replacements: {result.get('replacements_count')}")
        
        # Verify only 1 replacement made
        assert result.get('replacements_count') == 1, f"Expected 1 replacement, got {result.get('replacements_count')}"
        
        print("\n✅ Test passed!")


async def test_specific_slide():
    """Test replacement on specific slide only."""
    print("\n" + "="*60)
    print("Test: Specific slide replacement")
    print("="*60)
    
    with tempfile.TemporaryDirectory() as tmpdir:
        test_file = Path(tmpdir) / "test.pptx"
        create_test_pptx(test_file)
        
        # Perform replacement on slide 1 only - replace "note" with "document"
        result = await handle_replace_text({
            "pptx_path": str(test_file),
            "target": "slide_notes",
            "pattern": "note",
            "replacement": "document",
            "use_regex": False,
            "slide_number": 1,
            "in_place": True,
        })
        
        print(f"Success: {result.get('success')}")
        print(f"Slides scanned: {result.get('slides_scanned')}")
        print(f"Slides changed: {result.get('slides_changed')}")
        print(f"Affected slides: {result.get('affected_slides')}")
        
        # Verify only slide 1 was modified
        prs = Presentation(str(test_file))
        slide1_notes = prs.slides[0].notes_slide.notes_text_frame.text
        slide2_notes = prs.slides[1].notes_slide.notes_text_frame.text
        
        print(f"\nSlide 1 notes: {slide1_notes}")
        print(f"Slide 2 notes: {slide2_notes}")
        
        assert "document" in slide1_notes, "Slide 1 not modified"
        assert "note" not in slide1_notes, "Original text still in slide 1"
        assert "notes" in slide2_notes, "Slide 2 should be unchanged"
        assert "document" not in slide2_notes, "Slide 2 was incorrectly modified"
        
        print("\n✅ Test passed!")


async def test_output_path():
    """Test output to different file (not in-place)."""
    print("\n" + "="*60)
    print("Test: Output to different file")
    print("="*60)
    
    with tempfile.TemporaryDirectory() as tmpdir:
        test_file = Path(tmpdir) / "test.pptx"
        output_file = Path(tmpdir) / "output.pptx"
        create_test_pptx(test_file)
        
        # Get original content
        prs_before = Presentation(str(test_file))
        notes_before = prs_before.slides[0].notes_slide.notes_text_frame.text
        
        # Perform replacement to new file
        result = await handle_replace_text({
            "pptx_path": str(test_file),
            "target": "slide_notes",
            "pattern": "test",
            "replacement": "sample",
            "use_regex": False,
            "in_place": False,
            "output_path": str(output_file),
        })
        
        print(f"Success: {result.get('success')}")
        print(f"In place: {result.get('in_place')}")
        print(f"Output path: {result.get('pptx_path')}")
        
        # Verify original file unchanged
        prs_original = Presentation(str(test_file))
        notes_original = prs_original.slides[0].notes_slide.notes_text_frame.text
        
        assert notes_before == notes_original, "Original file was modified!"
        
        # Verify output file was created and modified
        assert output_file.exists(), "Output file not created!"
        prs_output = Presentation(str(output_file))
        notes_output = prs_output.slides[0].notes_slide.notes_text_frame.text
        
        assert "sample" in notes_output, "Output file not modified correctly"
        
        print("\n✅ Test passed!")


async def run_all_tests():
    """Run all tests."""
    print("\n" + "="*60)
    print("Running replace_text Tool Tests")
    print("="*60)
    
    tests = [
        ("Literal replacement in notes", test_replace_in_notes_literal),
        ("Literal replacement in content", test_replace_in_content_literal),
        ("Regex replacement with capture groups", test_replace_with_regex),
        ("Dry run mode", test_dry_run),
        ("Max replacements limit", test_max_replacements),
        ("Specific slide replacement", test_specific_slide),
        ("Output to different file", test_output_path),
    ]
    
    passed = 0
    failed = 0
    
    for test_name, test_func in tests:
        try:
            await test_func()
            passed += 1
        except Exception as e:
            failed += 1
            print(f"\n❌ Test failed: {test_name}")
            print(f"Error: {e}")
            import traceback
            traceback.print_exc()
    
    print("\n" + "="*60)
    print(f"Test Results: {passed} passed, {failed} failed")
    print("="*60)
    
    return failed == 0


if __name__ == "__main__":
    import asyncio
    success = asyncio.run(run_all_tests())
    exit(0 if success else 1)
