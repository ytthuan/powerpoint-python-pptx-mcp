"""Integration tests for slide visibility MCP tools."""

import asyncio
import tempfile
import os
from pptx import Presentation

from mcp_server.tools.read_tools import (
    handle_read_slide_content,
    handle_read_slides_metadata,
    handle_read_presentation_info,
)
from mcp_server.tools.slide_tools import handle_set_slide_visibility


def create_temp_pptx():
    """Create a temporary PPTX file and return its path."""
    fd, path = tempfile.mkstemp(suffix='.pptx')
    os.close(fd)  # Close the file descriptor
    return path


async def test_read_slides_metadata_tool():
    """Test read_slides_metadata MCP tool."""
    # Create a test presentation
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3 = prs.slides.add_slide(prs.slide_layouts[0])
    
    slide1.shapes.title.text = "Visible Slide 1"
    slide2.shapes.title.text = "Hidden Slide 2"
    slide3.shapes.title.text = "Visible Slide 3"
    
    # Hide slide 2
    sldIdLst = prs.slides._sldIdLst
    sldIdLst[1].set('show', '0')
    
    temp_file = create_temp_pptx()
    prs.save(temp_file)
    
    try:
        # Test with include_hidden=True
        result = await handle_read_slides_metadata({
            "pptx_path": temp_file,
            "include_hidden": True
        })
        
        assert result['total_slides'] == 3
        assert len(result['slides']) == 3
        assert result['slides'][0]['title'] == "Visible Slide 1"
        assert result['slides'][0]['hidden'] == False
        assert result['slides'][1]['title'] == "Hidden Slide 2"
        assert result['slides'][1]['hidden'] == True
        assert result['slides'][2]['title'] == "Visible Slide 3"
        assert result['slides'][2]['hidden'] == False
        
        # Test with include_hidden=False
        result = await handle_read_slides_metadata({
            "pptx_path": temp_file,
            "include_hidden": False
        })
        
        assert result['total_slides'] == 3
        assert len(result['slides']) == 2  # Only visible slides
        assert result['slides'][0]['title'] == "Visible Slide 1"
        assert result['slides'][1]['title'] == "Visible Slide 3"
        
        print("✓ test_read_slides_metadata_tool passed")
    finally:
        os.unlink(temp_file)


async def test_read_slide_content_with_hidden():
    """Test read_slide_content includes hidden status."""
    # Create a test presentation
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    
    slide1.shapes.title.text = "Visible Slide"
    slide2.shapes.title.text = "Hidden Slide"
    
    # Hide slide 2
    sldIdLst = prs.slides._sldIdLst
    sldIdLst[1].set('show', '0')
    
    temp_file = create_temp_pptx()
    prs.save(temp_file)
    
    try:
        # Test single slide
        result = await handle_read_slide_content({
            "pptx_path": temp_file,
            "slide_number": 1
        })
        
        assert 'hidden' in result
        assert result['hidden'] == False
        assert result['title'] == "Visible Slide"
        
        result = await handle_read_slide_content({
            "pptx_path": temp_file,
            "slide_number": 2
        })
        
        assert result['hidden'] == True
        assert result['title'] == "Hidden Slide"
        
        # Test all slides with include_hidden=True
        result = await handle_read_slide_content({
            "pptx_path": temp_file,
            "include_hidden": True
        })
        
        assert len(result['slides']) == 2
        assert result['slides'][0]['hidden'] == False
        assert result['slides'][1]['hidden'] == True
        
        # Test all slides with include_hidden=False
        result = await handle_read_slide_content({
            "pptx_path": temp_file,
            "include_hidden": False
        })
        
        assert len(result['slides']) == 1  # Only visible
        assert result['slides'][0]['hidden'] == False
        assert result['slides'][0]['title'] == "Visible Slide"
        
        print("✓ test_read_slide_content_with_hidden passed")
    finally:
        os.unlink(temp_file)


async def test_read_presentation_info_with_visibility():
    """Test read_presentation_info includes visibility stats."""
    # Create a test presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.slides.add_slide(prs.slide_layouts[1])
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.slides.add_slide(prs.slide_layouts[1])
    prs.slides.add_slide(prs.slide_layouts[0])
    
    # Hide slides 2, 3, and 5
    sldIdLst = prs.slides._sldIdLst
    sldIdLst[1].set('show', '0')
    sldIdLst[2].set('show', '0')
    sldIdLst[4].set('show', '0')
    
    temp_file = create_temp_pptx()
    prs.save(temp_file)
    
    try:
        result = await handle_read_presentation_info({
            "pptx_path": temp_file
        })
        
        assert 'visible_slides' in result
        assert 'hidden_slides' in result
        assert result['slide_count'] == 5
        assert result['visible_slides'] == 2
        assert result['hidden_slides'] == 3
        
        print("✓ test_read_presentation_info_with_visibility passed")
    finally:
        os.unlink(temp_file)


async def test_set_slide_visibility_tool():
    """Test set_slide_visibility MCP tool."""
    # Create a test presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.slides.add_slide(prs.slide_layouts[1])
    
    temp_file = create_temp_pptx()
    prs.save(temp_file)
    
    try:
        # Hide slide 1
        result = await handle_set_slide_visibility({
            "pptx_path": temp_file,
            "slide_number": 1,
            "hidden": True
        })
        
        assert result['success'] == True
        assert result['slide_number'] == 1
        assert result['hidden'] == True
        output_path = result['output_path']
        
        # Verify the slide is hidden in the output file
        verify_result = await handle_read_slide_content({
            "pptx_path": output_path,
            "slide_number": 1
        })
        assert verify_result['hidden'] == True
        
        # Show the slide again
        result = await handle_set_slide_visibility({
            "pptx_path": output_path,
            "slide_number": 1,
            "hidden": False
        })
        
        assert result['success'] == True
        assert result['hidden'] == False
        output_path2 = result['output_path']
        
        # Verify the slide is now visible
        verify_result = await handle_read_slide_content({
            "pptx_path": output_path2,
            "slide_number": 1
        })
        assert verify_result['hidden'] == False
        
        # Clean up output files
        if os.path.exists(output_path):
            os.unlink(output_path)
        if os.path.exists(output_path2):
            os.unlink(output_path2)
        
        print("✓ test_set_slide_visibility_tool passed")
    finally:
        os.unlink(temp_file)


async def main():
    """Run all integration tests."""
    print("Running slide visibility integration tests...\n")
    
    await test_read_slides_metadata_tool()
    await test_read_slide_content_with_hidden()
    await test_read_presentation_info_with_visibility()
    await test_set_slide_visibility_tool()
    
    print("\n✅ All integration tests passed!")


if __name__ == "__main__":
    asyncio.run(main())
