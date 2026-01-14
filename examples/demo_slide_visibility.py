"""Demo script for slide visibility features.

This script demonstrates:
1. Creating a presentation with hidden and visible slides
2. Reading slide metadata including hidden status
3. Filtering slides by visibility
4. Setting slide visibility
"""

import os
import tempfile

from pptx import Presentation

from mcp_server.core.pptx_handler import PPTXHandler


def demo_slide_visibility():
    """Demonstrate slide visibility features."""

    print("=" * 60)
    print("Slide Visibility Features Demo")
    print("=" * 60)

    # Create a test presentation
    print("\n1. Creating a test presentation with 5 slides...")
    prs = Presentation()

    for i in range(1, 6):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f"Slide {i}"

    # Hide slides 2 and 4
    print("   Hiding slides 2 and 4...")
    sldIdLst = prs.slides._sldIdLst
    sldIdLst[1].set("show", "0")  # Slide 2
    sldIdLst[3].set("show", "0")  # Slide 4

    # Save to temp file
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
        temp_path = temp_file.name
    prs.save(temp_path)
    print(f"   Saved to: {temp_path}")

    # Create handler
    handler = PPTXHandler(temp_path)

    # 2. Read presentation info
    print("\n2. Reading presentation info:")
    info = handler.get_presentation_info()
    print(f"   Total slides: {info['slide_count']}")
    print(f"   Visible slides: {info['visible_slides']}")
    print(f"   Hidden slides: {info['hidden_slides']}")

    # 3. Get all slides metadata
    print("\n3. Getting all slides metadata:")
    all_slides = handler.get_slides_metadata(include_hidden=True)
    for slide_meta in all_slides:
        status = "HIDDEN" if slide_meta["hidden"] else "VISIBLE"
        print(f"   Slide {slide_meta['slide_number']}: {slide_meta['title']} - {status}")

    # 4. Get only visible slides
    print("\n4. Getting only visible slides metadata:")
    visible_slides = handler.get_slides_metadata(include_hidden=False)
    for slide_meta in visible_slides:
        print(f"   Slide {slide_meta['slide_number']}: {slide_meta['title']}")

    # 5. Check individual slide visibility
    print("\n5. Checking individual slide visibility:")
    for i in range(1, 6):
        is_hidden = handler.is_slide_hidden(i)
        status = "hidden" if is_hidden else "visible"
        print(f"   Slide {i} is {status}")

    # 6. Change slide visibility
    print("\n6. Changing slide visibility:")
    print("   Making slide 2 visible...")
    handler.set_slide_hidden(2, False)
    print("   Making slide 3 hidden...")
    handler.set_slide_hidden(3, True)

    # 7. Verify changes
    print("\n7. Updated slide visibility:")
    for i in range(1, 6):
        is_hidden = handler.is_slide_hidden(i)
        status = "hidden" if is_hidden else "visible"
        print(f"   Slide {i} is now {status}")

    # 8. Get slide content with hidden status
    print("\n8. Reading slide content (includes hidden status):")
    content = handler.get_slide_content(2)
    print(f"   Slide 2: '{content['title']}' - hidden={content['hidden']}")

    content = handler.get_slide_content(3)
    print(f"   Slide 3: '{content['title']}' - hidden={content['hidden']}")

    # 9. Save changes
    print("\n9. Saving changes...")
    with tempfile.NamedTemporaryFile(suffix="_updated.pptx", delete=False) as output_file:
        output_path = output_file.name
    handler.save(output_path)
    print(f"   Saved updated presentation to: {output_path}")

    # 10. Verify persistence
    print("\n10. Verifying persistence after reload:")
    handler2 = PPTXHandler(output_path)
    info2 = handler2.get_presentation_info()
    print(f"    Total slides: {info2['slide_count']}")
    print(f"    Visible slides: {info2['visible_slides']}")
    print(f"    Hidden slides: {info2['hidden_slides']}")

    # Clean up
    os.unlink(temp_path)
    os.unlink(output_path)

    print("\n" + "=" * 60)
    print("Demo completed successfully!")
    print("=" * 60)

    print("\n📚 Use Cases:")
    print("  • Filter slides for presentation mode")
    print("  • Create speaker notes with backup slides")
    print("  • Build presentations with optional content")
    print("  • Manage slide visibility programmatically")


if __name__ == "__main__":
    demo_slide_visibility()
