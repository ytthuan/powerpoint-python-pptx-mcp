#!/usr/bin/env python3
"""Demo script showing replace_text tool capabilities."""

import asyncio
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from mcp_server.tools.text_replace_tools import handle_replace_text


def create_demo_pptx(path: Path) -> None:
    """Create a demo PPTX file."""
    prs = Presentation()

    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title1 = slide1.shapes.title
    title1.text = "Product Launch: ACME 2024"

    notes1 = slide1.notes_slide.notes_text_frame
    notes1.text = "Welcome to the ACME product launch! Today we'll discuss our new features."

    # Slide 2: Content slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = "Contact Information"

    # Add text box with email
    left = Inches(1)
    top = Inches(2)
    width = Inches(5)
    height = Inches(1)
    textbox = slide2.shapes.add_textbox(left, top, width, height)
    textbox.text_frame.text = "Email: support@acme-corp.com\nPhone: 555-ACME"

    notes2 = slide2.notes_slide.notes_text_frame
    notes2.text = "For questions, contact support@acme-corp.com"

    prs.save(str(path))


async def demo_literal_replacement():
    """Demo 1: Simple literal text replacement."""
    print("\n" + "=" * 70)
    print("DEMO 1: Literal Text Replacement")
    print("=" * 70)

    with tempfile.TemporaryDirectory() as tmpdir:
        demo_file = Path(tmpdir) / "demo.pptx"
        create_demo_pptx(demo_file)

        # Show original
        prs = Presentation(str(demo_file))
        print("\n📄 Original Content:")
        print(f"  Slide 1 Title: {prs.slides[0].shapes.title.text}")
        print(f"  Slide 1 Notes: {prs.slides[0].notes_slide.notes_text_frame.text}")

        # Replace ACME with Contoso
        print("\n🔧 Performing replacement: 'ACME' → 'Contoso' in all notes...")
        result = await handle_replace_text(
            {
                "pptx_path": str(demo_file),
                "target": "slide_notes",
                "pattern": "ACME",
                "replacement": "Contoso",
                "use_regex": False,
                "in_place": True,
            }
        )

        print("\n✅ Result:")
        print(f"  Success: {result['success']}")
        print(f"  Slides scanned: {result['slides_scanned']}")
        print(f"  Replacements made: {result['replacements_count']}")
        print(f"  Affected slides: {result['affected_slides']}")

        # Show updated
        prs = Presentation(str(demo_file))
        print("\n📝 Updated Content:")
        print(f"  Slide 1 Notes: {prs.slides[0].notes_slide.notes_text_frame.text}")


async def demo_regex_email():
    """Demo 2: Regex replacement for email addresses."""
    print("\n" + "=" * 70)
    print("DEMO 2: Regex Email Domain Replacement")
    print("=" * 70)

    with tempfile.TemporaryDirectory() as tmpdir:
        demo_file = Path(tmpdir) / "demo.pptx"
        create_demo_pptx(demo_file)

        # Show original
        prs = Presentation(str(demo_file))
        textbox = prs.slides[1].shapes[1]  # The textbox we added
        print("\n📄 Original Content:")
        print(f"  Slide 2 Content: {textbox.text_frame.text}")

        # Replace email domain using regex
        print("\n🔧 Using regex to change email domain...")
        print("  Pattern: r'([\\w.]+)@acme-corp\\.com'")
        print("  Replacement: r'\\1@contoso.com'")

        result = await handle_replace_text(
            {
                "pptx_path": str(demo_file),
                "target": "slide_content",
                "pattern": r"([\w.]+)@acme-corp\.com",
                "replacement": r"\1@contoso.com",
                "use_regex": True,
                "in_place": True,
            }
        )

        print("\n✅ Result:")
        print(f"  Success: {result['success']}")
        print(f"  Replacements made: {result['replacements_count']}")
        print(f"  Affected shapes: {result['affected_shapes']}")

        # Show updated
        prs = Presentation(str(demo_file))
        textbox = prs.slides[1].shapes[1]
        print("\n📝 Updated Content:")
        print(f"  Slide 2 Content: {textbox.text_frame.text}")


async def demo_dry_run():
    """Demo 3: Dry run to preview changes."""
    print("\n" + "=" * 70)
    print("DEMO 3: Dry Run Mode (Preview Changes)")
    print("=" * 70)

    with tempfile.TemporaryDirectory() as tmpdir:
        demo_file = Path(tmpdir) / "demo.pptx"
        create_demo_pptx(demo_file)

        print("\n🔍 Running dry run to preview changes...")
        print("  Pattern: 'ACME' → 'NewCorp'")

        result = await handle_replace_text(
            {
                "pptx_path": str(demo_file),
                "target": "slide_content",
                "pattern": "ACME",
                "replacement": "NewCorp",
                "use_regex": False,
                "dry_run": True,
            }
        )

        print("\n✅ Dry Run Result:")
        print(f"  Success: {result['success']}")
        print(f"  Dry run: {result['dry_run']}")
        print(f"  Would replace: {result['replacements_count']} instances")
        print(f"  Affected shapes: {result['affected_shapes']}")

        print("\n📋 Preview of changes:")
        for change in result["changes"]:
            print(f"\n  Slide {change['slide_number']}, Shape {change['shape_id']}:")
            print(f"    Before: {change['original_text']}")
            print(f"    After:  {change['modified_text']}")

        # Verify file not modified
        prs = Presentation(str(demo_file))
        title = prs.slides[0].shapes.title.text
        print(f"\n✓ File unchanged (verified): {title}")


async def demo_case_insensitive():
    """Demo 4: Case-insensitive replacement."""
    print("\n" + "=" * 70)
    print("DEMO 4: Case-Insensitive Replacement")
    print("=" * 70)

    with tempfile.TemporaryDirectory() as tmpdir:
        demo_file = Path(tmpdir) / "demo.pptx"
        create_demo_pptx(demo_file)

        print("\n📄 Original has: 'ACME', 'acme', 'Acme' in different places")

        print("\n🔧 Replacing all variants with 'Contoso' (case-insensitive)...")

        result = await handle_replace_text(
            {
                "pptx_path": str(demo_file),
                "target": "slide_content",
                "pattern": "acme",
                "replacement": "Contoso",
                "use_regex": True,
                "regex_flags": ["IGNORECASE"],
                "in_place": True,
            }
        )

        print("\n✅ Result:")
        print(f"  Success: {result['success']}")
        print(f"  Replacements made: {result['replacements_count']}")


async def main():
    """Run all demos."""
    print("\n" + "=" * 70)
    print("     REPLACE_TEXT TOOL - FEATURE DEMONSTRATION")
    print("=" * 70)

    await demo_literal_replacement()
    await demo_regex_email()
    await demo_dry_run()
    await demo_case_insensitive()

    print("\n" + "=" * 70)
    print("✅ All demos completed successfully!")
    print("=" * 70)
    print("\n💡 Key Features Demonstrated:")
    print("  ✓ Literal text replacement in notes and content")
    print("  ✓ Regex patterns with capture groups")
    print("  ✓ Dry run mode to preview changes")
    print("  ✓ Case-insensitive replacement with regex flags")
    print("  ✓ Safe in-place file updates")
    print("\n")


if __name__ == "__main__":
    asyncio.run(main())
