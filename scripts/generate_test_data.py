from pptx import Presentation
from pathlib import Path


def generate_test_presentation():
    prs = Presentation()
    for i in range(1, 6):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f"Slide {i} Title"
        # Add some text to the body
        if slide.placeholders:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:  # body
                    ph.text = f"This is the body text for slide {i}."

        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = f"Speaker notes for slide {i}."

    output_path = Path("tests/integration/test_data/test_presentation.pptx")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Created {output_path}")


if __name__ == "__main__":
    generate_test_presentation()
