#!/usr/bin/env python3
"""Export and update PowerPoint speaker notes (Notes Slide) in a .pptx.

⚠️ DEPRECATION WARNING ⚠️
==========================
This script is DEPRECATED and will be removed in a future version.
Please use the MCP server tools instead for all PPTX operations.

Migration guide: See MIGRATION.md for instructions on transitioning to MCP tools.

The MCP server provides:
- Better error handling and validation
- Consistent API across all operations
- Enhanced security features
- Structured logging and monitoring
- Type safety and comprehensive testing

For MCP server usage, see README.md and AGENTS.md documentation.
==========================

Uses python-pptx: https://python-pptx.readthedocs.io/en/latest/user/notes.html

Design goal: only modify slide notes; do not alter slide shapes/content.
"""

from __future__ import annotations

import argparse
import datetime as _dt
import json
import os
import sys
import tempfile
import warnings
import zipfile
from pathlib import Path
from typing import Any, Dict, Iterable, List, Tuple

from lxml import etree
from pptx import Presentation


# Show deprecation warning
def _show_deprecation_warning():
    """Display deprecation warning to users."""
    warning_msg = """
╔═══════════════════════════════════════════════════════════════════════╗
║                      ⚠️  DEPRECATION WARNING ⚠️                        ║
╠═══════════════════════════════════════════════════════════════════════╣
║  This script (pptx_notes.py) is DEPRECATED and will be removed       ║
║  in a future version.                                                  ║
║                                                                        ║
║  Please use the MCP server tools instead:                             ║
║  - Better error handling and validation                               ║
║  - Consistent API across all operations                               ║
║  - Enhanced security features                                         ║
║  - Structured logging and monitoring                                  ║
║                                                                        ║
║  See MIGRATION.md for migration instructions.                         ║
║  See README.md for MCP server usage.                                  ║
╚═══════════════════════════════════════════════════════════════════════╝
"""
    print(warning_msg, file=sys.stderr)
    warnings.warn(
        "pptx_notes.py is deprecated. Use MCP server tools instead. "
        "See MIGRATION.md for migration guide.",
        DeprecationWarning,
        stacklevel=2
    )


def _slide_title(slide) -> str:
    title_shape = getattr(slide.shapes, "title", None)
    if not title_shape:
        return ""
    try:
        return title_shape.text_frame.text.strip()
    except Exception:
        # Some layouts may not have a standard title placeholder.
        return ""


def _read_notes_text(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    try:
        return slide.notes_slide.notes_text_frame.text
    except Exception:
        # If the notes placeholder is missing/corrupt, fail soft.
        return ""


def dump_notes(pptx_path: Path) -> Dict[str, Any]:
    pres = Presentation(str(pptx_path))

    slides: List[Dict[str, Any]] = []
    for idx, slide in enumerate(pres.slides, start=1):
        slides.append(
            {
                "slide": idx,
                "title": _slide_title(slide),
                "notes": _read_notes_text(slide),
            }
        )

    return {
        "source": pptx_path.name,
        "generated_at": _dt.datetime.now(tz=_dt.timezone.utc).isoformat(),
        "slides": slides,
    }


def _iter_updates(payload: Any) -> Iterable[Tuple[int, str]]:
    """Yield (slide_number, notes_text) updates from supported JSON shapes.

    Supported formats:
    1) {"slides": [{"slide": 1, "notes": "..."}, ...]}
    2) {"1": "...", "2": "..."}  (mapping of slide->notes)
    """

    if isinstance(payload, dict) and "slides" in payload and isinstance(payload["slides"], list):
        for item in payload["slides"]:
            if not isinstance(item, dict):
                raise ValueError("Invalid JSON: slides[] items must be objects")
            slide_no = item.get("slide")
            notes = item.get("notes")
            if not isinstance(slide_no, int):
                raise ValueError("Invalid JSON: slides[] item missing integer 'slide'")
            if notes is None:
                notes = ""
            if not isinstance(notes, str):
                raise ValueError("Invalid JSON: slides[] item 'notes' must be string")
            yield slide_no, notes
        return

    if isinstance(payload, dict):
        for k, v in payload.items():
            try:
                slide_no = int(k)
            except Exception as e:
                raise ValueError(f"Invalid JSON mapping key (expected slide number): {k!r}") from e
            if v is None:
                v = ""
            if not isinstance(v, str):
                raise ValueError(f"Invalid JSON mapping value for slide {k!r}: must be string")
            yield slide_no, v
        return

    raise ValueError("Invalid JSON: expected object with 'slides' list or mapping of slide->notes")


def apply_notes(pptx_in: Path, updates_json: Path, pptx_out: Path) -> None:
    pres = Presentation(str(pptx_in))

    payload = json.loads(updates_json.read_text(encoding="utf-8"))
    updates = list(_iter_updates(payload))

    changed = 0
    for slide_no, notes_text in updates:
        if slide_no < 1 or slide_no > len(pres.slides):
            raise ValueError(f"Slide {slide_no} is out of range (1..{len(pres.slides)})")
        slide = pres.slides[slide_no - 1]
        # Accessing .notes_slide creates notes slide if it doesn't exist.
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text
        changed += 1

    pres.save(str(pptx_out))

    print(f"Updated notes for {changed} slide(s).")
    print(f"Wrote: {pptx_out}")


def apply_notes_in_place(pptx_in: Path, updates_json: Path) -> None:
    """Safely overwrite the input PPTX by writing to a temp file then replacing."""

    pptx_in = pptx_in.resolve()
    tmp_dir = str(pptx_in.parent)

    fd, tmp_path = tempfile.mkstemp(prefix=pptx_in.stem + ".", suffix=".tmp.pptx", dir=tmp_dir)
    os.close(fd)
    tmp_file = Path(tmp_path)
    try:
        apply_notes(pptx_in, updates_json, tmp_file)
        os.replace(str(tmp_file), str(pptx_in))
        print(f"Overwrote: {pptx_in}")
    finally:
        if tmp_file.exists():
            try:
                tmp_file.unlink()
            except Exception:
                pass


_REL_NS = {
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}

_XML_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}

# For notes slides, txBody is in 'p' namespace, not 'a'
_NOTES_XML_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def _notes_part_for_slide(zip_in: zipfile.ZipFile, slide_no: int) -> str | None:
    rel_path = f"ppt/slides/_rels/slide{slide_no}.xml.rels"
    try:
        rel_xml = zip_in.read(rel_path)
    except KeyError:
        return None

    root = etree.fromstring(rel_xml)
    for rel in root.findall("rel:Relationship", namespaces=_REL_NS):
        if rel.get("Type") == "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide":
            target = rel.get("Target")
            if not target:
                return None
            # Typically "../notesSlides/notesSlideX.xml". Normalize to zip path.
            target = target.replace("\\", "/")
            if target.startswith("../"):
                target = target[len("../") :]
            return f"ppt/{target}"
    return None


def _set_notes_text(notes_xml: bytes, notes_text: str) -> bytes:
    """Replace the text content of the notes placeholder while preserving everything else."""

    parser = etree.XMLParser(remove_blank_text=False)
    root = etree.fromstring(notes_xml, parser)

    # Find the 'body' placeholder shape which holds the notes text.
    # Use full XPath because ElementPath.find() has limited predicate support.
    sp_list = root.xpath("//p:sp[p:nvSpPr/p:nvPr/p:ph[@type='body']]", namespaces=_XML_NS)
    sp = sp_list[0] if sp_list else None
    if sp is None:
        # Fallback: first shape with a text body.
        sp_list = root.xpath("//p:sp[.//a:txBody]", namespaces=_XML_NS)
        sp = sp_list[0] if sp_list else None
    if sp is None:
        return notes_xml

    # Try both namespaces - txBody can be in 'a' or 'p' namespace depending on structure
    # In notes slides, txBody is in 'p' namespace
    tx_list = sp.xpath(".//p:txBody", namespaces=_XML_NS)
    if not tx_list:
        # Fallback: try 'a' namespace (for regular slides)
        tx_list = sp.xpath(".//a:txBody", namespaces=_XML_NS)
    tx_body = tx_list[0] if tx_list else None
    if tx_body is None:
        return notes_xml

    # Keep bodyPr and lstStyle, replace only paragraphs.
    # Paragraphs are in 'a' namespace even when txBody is in 'p' namespace
    for p_el in list(tx_body.findall("a:p", namespaces=_XML_NS)):
        tx_body.remove(p_el)

    # Clean the text: replace control characters except newlines
    # Replace \u000b (vertical tab) and other control chars with newline or space
    cleaned_text = notes_text.replace("\r\n", "\n").replace("\r", "\n")
    # Replace vertical tab and other control characters with newline
    import re
    cleaned_text = re.sub(r'[\x00-\x08\x0b-\x0c\x0e-\x1f]', '\n', cleaned_text)
    
    lines = cleaned_text.split("\n")
    if len(lines) == 0:
        lines = [""]

    for line in lines:
        p_el = etree.Element(f"{{{_XML_NS['a']}}}p")
        r_el = etree.SubElement(p_el, f"{{{_XML_NS['a']}}}r")
        
        # Check if this line should be bold (starts with "- short version:" or "- original:")
        line_stripped = line.strip()
        is_bold = line_stripped.startswith("- Short version:") or line_stripped.startswith("- Original:")
        
        if is_bold:
            # Add run properties with bold formatting
            r_pr = etree.SubElement(r_el, f"{{{_XML_NS['a']}}}rPr")
            r_pr.set("b", "1")  # Bold
        
        t_el = etree.SubElement(r_el, f"{{{_XML_NS['a']}}}t")
        t_el.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
        # Remove any remaining control characters that XML doesn't allow
        clean_line = ''.join(char for char in line if ord(char) >= 32 or char in '\n\r\t')
        t_el.text = clean_line
        tx_body.append(p_el)

    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=None)


def apply_notes_zip(pptx_in: Path, updates_json: Path, pptx_out: Path) -> None:
    """Update notes by editing only notesSlide XML parts inside the PPTX zip."""

    payload = json.loads(updates_json.read_text(encoding="utf-8"))
    updates = list(_iter_updates(payload))

    with zipfile.ZipFile(pptx_in, "r") as zin:
        # Copy all parts; overwrite only the notes slide parts we touch.
        notes_updates: Dict[str, bytes] = {}
        for slide_no, notes_text in updates:
            notes_part = _notes_part_for_slide(zin, slide_no)
            if not notes_part:
                # No notes part exists; skip unless caller wants creation (not supported).
                continue
            original_notes_xml = zin.read(notes_part)
            notes_updates[notes_part] = _set_notes_text(original_notes_xml, notes_text)

        with zipfile.ZipFile(pptx_out, "w", compression=zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename in notes_updates:
                    data = notes_updates[item.filename]
                zi = zipfile.ZipInfo(filename=item.filename, date_time=item.date_time)
                zi.compress_type = zipfile.ZIP_DEFLATED
                zi.external_attr = item.external_attr
                zi.comment = item.comment
                zi.extra = item.extra
                zi.internal_attr = item.internal_attr
                zi.create_system = item.create_system
                zout.writestr(zi, data)

    print(f"Updated notes for {len(updates)} slide(s) (zip mode).")
    print(f"Wrote: {pptx_out}")


def apply_notes_zip_in_place(pptx_in: Path, updates_json: Path) -> None:
    """Safely overwrite PPTX using zip-mode notes editing."""

    pptx_in = pptx_in.resolve()
    tmp_dir = str(pptx_in.parent)
    fd, tmp_path = tempfile.mkstemp(prefix=pptx_in.stem + ".", suffix=".tmp.pptx", dir=tmp_dir)
    os.close(fd)
    tmp_file = Path(tmp_path)
    try:
        apply_notes_zip(pptx_in, updates_json, tmp_file)
        os.replace(str(tmp_file), str(pptx_in))
        print(f"Overwrote: {pptx_in}")
    finally:
        if tmp_file.exists():
            try:
                tmp_file.unlink()
            except Exception:
                pass


def _default_out_path(pptx_in: Path) -> Path:
    # Avoid in-place writes by default.
    return pptx_in.with_name(pptx_in.stem + ".notes.pptx")


def main(argv: List[str] | None = None) -> int:
    """Main entry point for pptx_notes script.
    
    ⚠️ DEPRECATED: Use MCP server tools instead. See MIGRATION.md
    """
    # Show deprecation warning
    _show_deprecation_warning()
    
    parser = argparse.ArgumentParser(
        prog="pptx_notes",
        description="Read/update PowerPoint slide speaker notes (.pptx) using python-pptx. "
                    "⚠️ DEPRECATED - Use MCP server tools instead.",
    )
    sub = parser.add_subparsers(dest="cmd", required=True)

    p_dump = sub.add_parser("dump", help="Export notes to JSON")
    p_dump.add_argument("pptx", type=Path, help="Input .pptx path")
    p_dump.add_argument(
        "--out",
        type=Path,
        default=None,
        help="Output JSON path (default: <pptx>.notes.json)",
    )

    p_apply = sub.add_parser("apply", help="Apply notes from JSON and write a new .pptx")
    p_apply.add_argument("pptx", type=Path, help="Input .pptx path")
    p_apply.add_argument("updates", type=Path, help="JSON file produced by dump (or mapping slide->notes)")
    p_apply.add_argument(
        "--out",
        type=Path,
        default=None,
        help="Output .pptx path (default: <pptx>.notes.pptx)",
    )
    p_apply.add_argument(
        "--in-place",
        action="store_true",
        help="Overwrite the input PPTX (writes to temp file then replaces)",
    )
    p_apply.add_argument(
        "--engine",
        choices=["zip", "python-pptx"],
        default="zip",
        help="How to write notes: 'zip' edits only notes parts (safest); 'python-pptx' rewrites the file.",
    )

    args = parser.parse_args(argv)

    if args.cmd == "dump":
        pptx_path: Path = args.pptx
        if not pptx_path.exists():
            parser.error(f"Input not found: {pptx_path}")
        out: Path = args.out or pptx_path.with_name(pptx_path.stem + ".notes.json")
        data = dump_notes(pptx_path)
        out.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        print(f"Wrote: {out}")
        return 0

    if args.cmd == "apply":
        pptx_in: Path = args.pptx
        updates_json: Path = args.updates
        if not pptx_in.exists():
            parser.error(f"Input not found: {pptx_in}")
        if not updates_json.exists():
            parser.error(f"Updates JSON not found: {updates_json}")
        if args.in_place and args.out is not None:
            parser.error("Use either --in-place or --out, not both")

        if args.engine == "python-pptx":
            if args.in_place:
                apply_notes_in_place(pptx_in, updates_json)
            else:
                pptx_out: Path = args.out or _default_out_path(pptx_in)
                apply_notes(pptx_in, updates_json, pptx_out)
        else:
            if args.in_place:
                apply_notes_zip_in_place(pptx_in, updates_json)
            else:
                pptx_out: Path = args.out or _default_out_path(pptx_in)
                apply_notes_zip(pptx_in, updates_json, pptx_out)
        return 0

    parser.error("Unknown command")


if __name__ == "__main__":
    raise SystemExit(main())
